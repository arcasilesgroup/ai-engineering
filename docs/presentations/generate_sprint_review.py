"""Generate Sprint Review presentation (.pptx) - Feb 16 - Mar 16, 2026.

Output: docs/presentation/sprint-review-2026-03.pptx
Brand: ai-engineering dark mode — dark backgrounds, teal accents, technical tone.
Audience: engineers, heads, managers — technical but client-facing.
"""

from __future__ import annotations

import contextlib
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# ai-engineering Dark-Mode Colour Palette
# ---------------------------------------------------------------------------
AI_BG_DARK = RGBColor(0x0B, 0x11, 0x20)
AI_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
AI_TEXT_PRIMARY = RGBColor(0xE2, 0xE8, 0xF0)

AI_ACCENT = RGBColor(0x00, 0xD4, 0xAA)
AI_PRIMARY = RGBColor(0x1E, 0x3A, 0x5F)
AI_PRIMARY_LIGHT = RGBColor(0x2A, 0x4F, 0x7A)
AI_ERROR = RGBColor(0xEF, 0x44, 0x44)
AI_SUCCESS = RGBColor(0x10, 0xB9, 0x81)
AI_WARNING = RGBColor(0xF5, 0x9E, 0x0B)

AI_TEXT_LIGHT = RGBColor(0xF8, 0xFA, 0xFB)
AI_NEUTRAL = RGBColor(0x64, 0x74, 0x8B)
AI_TEXT_MUTED = RGBColor(0x94, 0xA3, 0xB8)

AI_BORDER_DARK = RGBColor(0x1A, 0x2A, 0x40)
AI_CARD_DARK = RGBColor(0x1E, 0x29, 0x3B)

SEC_BLUE = RGBColor(0x2E, 0x6B, 0xA4)
SEC_BLUE_LIGHT = RGBColor(0x7A, 0xB5, 0xD6)
SEC_TEAL = RGBColor(0x00, 0xD4, 0xAA)
SEC_TEAL_LIGHT = RGBColor(0x5C, 0xE8, 0xCC)
SEC_PURPLE = RGBColor(0x7B, 0x3F, 0xA0)
SEC_PURPLE_LIGHT = RGBColor(0xB3, 0x8B, 0xCF)

# ---------------------------------------------------------------------------
# Typography
# ---------------------------------------------------------------------------
FONT_TITLE = "JetBrains Mono"
FONT_BODY = "Inter"

# ---------------------------------------------------------------------------
# Slide dimensions (16:9)
# ---------------------------------------------------------------------------
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

LEFT_MARGIN = Inches(1.2)
CONTENT_W = SLIDE_W - LEFT_MARGIN - Inches(1.2)

_PT_3 = Pt(3)
_PT_14 = Pt(14)
_PT_20 = Pt(20)
_PT_1 = Pt(1)
_PT_12 = Pt(12)

# ---------------------------------------------------------------------------
# Speaker Notes
# ---------------------------------------------------------------------------
NOTES: dict[int, str] = {
    1: (
        "Welcome to the sprint review for the period February 16 through March 16, 2026.\n\n"
        "This has been an exceptionally productive sprint — 76 commits across 19 active days, "
        "touching over 1,100 files with 118,000 lines of new code.\n\n"
        "We delivered across architecture, observability, security, testing, and CI/CD."
    ),
    2: (
        "Let me give you a high-level overview.\n\n"
        "76 commits by the engineering team across 19 working days. "
        "Over 118,000 lines of code were added across 1,143 files. "
        "We completed or advanced 20+ specifications. "
        "Test coverage rose to 91%, well above our 80% target.\n\n"
        "The work spanned 6 major themes which we'll cover in the next slides."
    ),
    3: (
        "The biggest delivery this sprint is Architecture v3 — spec 051.\n\n"
        "This is a clean-sheet redesign of our governance surface. "
        "We went from 7 agents to 10, adding guard for proactive governance, "
        "guide for developer onboarding, and operate for SRE tasks.\n\n"
        "Skills expanded from 35 (5 were stubs) to 40, all fully functional. "
        "12 skills were renamed for self-documenting names — build became code, "
        "cicd became pipeline, db became schema.\n\n"
        "The self-improvement loop via the evolve skill analyzes audit-log data "
        "and proposes framework improvements automatically."
    ),
    4: (
        "Spec 053 addresses a critical behavioral gap.\n\n"
        "Previously, IDE integrations used thin wrappers that referenced "
        "agents and skills in .ai-engineering — but most IDEs don't execute "
        "these references, so behaviors like the Interrogation Phase never ran.\n\n"
        "The solution: full IDE-adapted copies in .claude, .github, and .agents "
        "directories. 9 migration phases, completed through Phase 9.\n\n"
        "Canonical source remains in src/ai_engineering/templates."
    ),
    5: (
        "Observability matured significantly this sprint.\n\n"
        "We now have 5 dashboard modes: engineer, team, AI, DORA, and health. "
        "Each surfaces different metrics — from code quality to deployment frequency.\n\n"
        "Cross-IDE telemetry via signals emit tracks skill invocations and agent dispatches. "
        "Rich formatting with progress bars, score badges, and color-coded metrics.\n\n"
        "Dual-output: human-readable Rich tables or JSON for automation."
    ),
    6: (
        "Security hardening was driven by SonarCloud Quality Gate findings.\n\n"
        "We fixed 5 vulnerabilities — 4 BLOCKER path traversal and 1 MAJOR command injection. "
        "All fixed by refactoring the code, not by suppressing findings.\n\n"
        "We added a no-suppression governance rule: no NOSONAR, nosec, noqa, or similar "
        "comments are allowed. Fix the root cause or escalate.\n\n"
        "Snyk integration provides continuous dependency vulnerability monitoring. "
        "Our 4-stage gate pipeline — pre-commit, commit-msg, pre-push, and CI — "
        "ensures zero ungated operations."
    ),
    7: (
        "Testing excellence through spec 052.\n\n"
        "The TDD protocol separates RED and GREEN phases — tests are written first "
        "and are immutable during implementation. The iron law: never weaken tests.\n\n"
        "The Interrogation Phase challenges assumptions before planning — "
        "mapping KNOWN, ASSUMED, and UNKNOWN.\n\n"
        "Coverage is at 91%. We support 20 stacks with consistent patterns: "
        "fakes over mocks, AAA pattern, dynamic counts.\n\n"
        "The test rewriting initiative moved from 63 individual patches per test "
        "to 1 comprehensive fixture — much more maintainable."
    ),
    8: (
        "CI/CD improvements focused on reliability and trust.\n\n"
        "Zero-rebuild release model: the release job now downloads the CI-validated "
        "artifact instead of rebuilding. What CI tests is exactly what ships.\n\n"
        "Dependabot PRs are now grouped and exempted from gate trailer checks "
        "to reduce noise while maintaining security.\n\n"
        "The CI matrix covers Linux, macOS, and Windows with Python 3.11+."
    ),
    9: (
        "Quality metrics at a glance.\n\n"
        "Coverage at 91% — target is 80%, so we're well above.\n"
        "Cyclomatic complexity capped at 10 per function, enforced by ruff.\n"
        "Cognitive complexity capped at 15.\n"
        "Gate pass rate is 100% — every commit passes all gates.\n"
        "Token efficiency at 99.19% — nearly all tokens are deferred, not preloaded.\n\n"
        "These aren't aspirational — they're enforced by our gate pipeline."
    ),
    10: (
        "The governance surface now stands at 10 agents, 40 skills, and 50+ standards.\n\n"
        "We've delivered 38 specifications total since inception. "
        "The framework supports 4 IDEs: Claude Code, GitHub Copilot, Gemini CLI, and Codex.\n\n"
        "Decision store has 17 recorded decisions with SHA-256 context hashing. "
        "Audit log tracks 185+ events for full traceability."
    ),
    11: (
        "Two active risks to watch.\n\n"
        "The ty type checker shows 84 diagnostics — these are non-blocking tech debt "
        "but should be addressed systematically.\n\n"
        "Cross-OS CI is stabilizing — Windows and macOS need further hardening.\n\n"
        "Next sprint priorities: complete spec-053 migration, "
        "release readiness with docs site (Nextra), PyPI stable release, "
        "and multi-IDE verification."
    ),
    12: (
        "Thank you for your time. Open for questions.\n\n"
        "Key takeaway: the framework evolved from content-only governance "
        "to a production-grade platform with measurable quality improvements "
        "across every dimension — security, testing, observability, and architecture."
    ),
}

# ---------------------------------------------------------------------------
# Helper utilities (same as board presentation)
# ---------------------------------------------------------------------------


def _font(
    run,
    *,
    name: str = FONT_BODY,
    size: Pt = _PT_14,
    color: RGBColor = AI_TEXT_LIGHT,
    bold: bool = False,
):
    run.font.name = name
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    font_name: str = FONT_BODY,
    font_size: Pt = _PT_14,
    color: RGBColor = AI_TEXT_LIGHT,
    bold: bool = False,
    alignment: PP_ALIGN = PP_ALIGN.LEFT,
    word_wrap: bool = True,
    anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = word_wrap
    tf.auto_size = None
    with contextlib.suppress(Exception):
        tf.paragraphs[0].alignment = alignment
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    _font(run, name=font_name, size=font_size, color=color, bold=bold)
    txbox.text_frame.paragraphs[0].space_before = Pt(0)
    txbox.text_frame.paragraphs[0].space_after = Pt(0)
    with contextlib.suppress(Exception):
        txbox.text_frame._txBody.bodyPr.set(
            "anchor",
            {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}.get(
                anchor, "t"
            ),
        )
    return txbox


def add_rich_textbox(
    slide, left, top, width, height, lines, *, alignment=PP_ALIGN.LEFT, line_spacing=_PT_20
):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = alignment
        if line_spacing:
            p.space_after = line_spacing
        run = p.add_run()
        run.text = ln.get("text", "")
        _font(
            run,
            name=ln.get("font_name", FONT_BODY),
            size=ln.get("font_size", Pt(14)),
            color=ln.get("color", AI_TEXT_LIGHT),
            bold=ln.get("bold", False),
        )
    return txbox


def add_accent_bar(slide, left, top, width, *, height=_PT_3, color=AI_ACCENT):
    bar = slide.shapes.add_shape(1, left, top, width, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar


def add_card(
    slide,
    left,
    top,
    width,
    height,
    *,
    fill_color=AI_CARD_DARK,
    border_color=None,
    border_width=_PT_1,
    left_accent_color=None,
    left_accent_width=None,
):
    card = slide.shapes.add_shape(1, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = fill_color
    if border_color:
        card.line.color.rgb = border_color
        card.line.width = border_width
    else:
        card.line.fill.background()
    if left_accent_color:
        _accent_w = left_accent_width if left_accent_width is not None else Inches(0.08)
        accent = slide.shapes.add_shape(1, left, top, _accent_w, height)
        accent.fill.solid()
        accent.fill.fore_color.rgb = left_accent_color
        accent.line.fill.background()
    return card


def add_rounded_rect(
    slide, left, top, width, height, *, fill_color=AI_CARD_DARK, border_color=AI_TEXT_MUTED
):
    rr = slide.shapes.add_shape(5, left, top, width, height)
    rr.fill.solid()
    rr.fill.fore_color.rgb = fill_color
    if border_color:
        rr.line.color.rgb = border_color
        rr.line.width = Pt(1)
    else:
        rr.line.fill.background()
    return rr


def add_rect(
    slide,
    left,
    top,
    width,
    height,
    *,
    fill_color=AI_CARD_DARK,
    border_color=AI_TEXT_MUTED,
    border_width=_PT_1,
):
    rect = slide.shapes.add_shape(1, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill_color
    if border_color:
        rect.line.color.rgb = border_color
        rect.line.width = border_width
    else:
        rect.line.fill.background()
    return rect


def set_notes(slide, text: str):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


def add_slide_header(slide, title: str, *, subtitle: str | None = None):
    add_accent_bar(slide, LEFT_MARGIN, Inches(0.6), CONTENT_W)
    add_textbox(
        slide,
        LEFT_MARGIN,
        Inches(0.75),
        CONTENT_W,
        Inches(0.7),
        text=title,
        font_name=FONT_TITLE,
        font_size=Pt(34),
        color=AI_WHITE,
        bold=True,
    )
    if subtitle:
        add_textbox(
            slide,
            LEFT_MARGIN,
            Inches(1.35),
            CONTENT_W,
            Inches(0.4),
            text=subtitle,
            font_name=FONT_BODY,
            font_size=Pt(18),
            color=AI_NEUTRAL,
        )


def _blank_slide(prs):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = AI_BG_DARK
    return slide


def add_styled_table(slide, left, top, width, height, rows, cols, *, header_fill=AI_ACCENT):
    shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    tbl = shape.table
    for ci in range(cols):
        cell = tbl.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
    return tbl, shape


def _style_cell(
    cell,
    text,
    *,
    font_size=_PT_12,
    bold=False,
    color=AI_TEXT_LIGHT,
    alignment=PP_ALIGN.LEFT,
    font_name=FONT_BODY,
    fill_color=None,
):
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    _font(run, name=font_name, size=font_size, color=color, bold=bold)
    cell.text_frame.word_wrap = True
    if fill_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill_color
    cell.margin_left = Inches(0.08)
    cell.margin_right = Inches(0.08)
    cell.margin_top = Inches(0.04)
    cell.margin_bottom = Inches(0.04)


# ---------------------------------------------------------------------------
# KPI card helper
# ---------------------------------------------------------------------------
def _add_kpi_card(slide, x, y, w, h, value, label, *, accent=AI_ACCENT):
    add_card(slide, x, y, w, h, fill_color=AI_CARD_DARK, border_color=AI_BORDER_DARK)
    add_textbox(
        slide,
        x,
        y + Inches(0.15),
        w,
        Inches(0.6),
        text=value,
        font_name=FONT_TITLE,
        font_size=Pt(36),
        color=accent,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        x,
        y + Inches(0.8),
        w,
        Inches(0.4),
        text=label,
        font_size=Pt(13),
        color=AI_NEUTRAL,
        alignment=PP_ALIGN.CENTER,
    )


# ---------------------------------------------------------------------------
# Feature card helper (for slides 3-8)
# ---------------------------------------------------------------------------
def _add_feature_card(slide, x, y, w, h, title, bullets, *, accent=AI_ACCENT):
    add_card(
        slide,
        x,
        y,
        w,
        h,
        fill_color=AI_CARD_DARK,
        border_color=AI_BORDER_DARK,
        left_accent_color=accent,
    )
    add_textbox(
        slide,
        x + Inches(0.2),
        y + Inches(0.1),
        w - Inches(0.4),
        Inches(0.35),
        text=title,
        font_size=Pt(14),
        color=AI_TEXT_PRIMARY,
        bold=True,
    )
    lines = [{"text": f"  {b}", "font_size": Pt(11), "color": AI_TEXT_MUTED} for b in bullets]
    add_rich_textbox(
        slide,
        x + Inches(0.2),
        y + Inches(0.45),
        w - Inches(0.4),
        h - Inches(0.55),
        lines,
        line_spacing=Pt(6),
    )


# ---------------------------------------------------------------------------
# Slide 1 — Title
# ---------------------------------------------------------------------------
def slide_01_title(prs):
    slide = _blank_slide(prs)

    add_accent_bar(slide, Inches(0), Inches(0.4), SLIDE_W, height=Pt(4))

    add_textbox(
        slide,
        LEFT_MARGIN,
        Inches(1.8),
        CONTENT_W,
        Inches(1.0),
        text="Sprint Review",
        font_name=FONT_TITLE,
        font_size=Pt(52),
        color=AI_WHITE,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        LEFT_MARGIN,
        Inches(2.85),
        CONTENT_W,
        Inches(0.6),
        text="ai-engineering",
        font_name=FONT_TITLE,
        font_size=Pt(28),
        color=AI_ACCENT,
        alignment=PP_ALIGN.CENTER,
    )

    bar_w = Inches(3)
    bar_left = (SLIDE_W - bar_w) // 2
    add_accent_bar(slide, bar_left, Inches(3.6), bar_w, height=Pt(3))

    add_textbox(
        slide,
        LEFT_MARGIN,
        Inches(3.9),
        CONTENT_W,
        Inches(0.5),
        text="February 16 — March 16, 2026",
        font_size=Pt(20),
        color=AI_TEXT_PRIMARY,
        alignment=PP_ALIGN.CENTER,
    )

    # Sprint badges
    badge_y = Inches(5.0)
    badge_w = Inches(2.0)
    badge_h = Inches(0.45)
    gap = Inches(0.3)
    badges = ["76 Commits", "20+ Specs", "118K+ LOC", "91% Coverage"]
    total_w = len(badges) * badge_w + (len(badges) - 1) * gap
    start_x = (SLIDE_W - total_w) // 2
    for i, label in enumerate(badges):
        x = start_x + i * (badge_w + gap)
        add_rounded_rect(
            slide, x, badge_y, badge_w, badge_h, fill_color=AI_CARD_DARK, border_color=AI_ACCENT
        )
        add_textbox(
            slide,
            x,
            badge_y,
            badge_w,
            badge_h,
            text=label,
            font_size=Pt(12),
            color=AI_ACCENT,
            bold=True,
            alignment=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )

    add_accent_bar(slide, Inches(0), Inches(7.0), SLIDE_W, height=Pt(4))
    set_notes(slide, NOTES[1])


# ---------------------------------------------------------------------------
# Slide 2 — Sprint Overview
# ---------------------------------------------------------------------------
def slide_02_overview(prs):
    slide = _blank_slide(prs)
    add_slide_header(slide, "Sprint Overview", subtitle="Feb 16 - Mar 16, 2026  |  19 active days")

    # KPI cards row
    kpi_y = Inches(2.2)
    kpi_w = Inches(2.4)
    kpi_h = Inches(1.3)
    kpi_gap = Inches(0.3)
    kpis = [
        ("76", "Commits"),
        ("118K+", "Lines Added"),
        ("20+", "Specs Delivered"),
        ("91%", "Test Coverage"),
    ]
    total = len(kpis) * kpi_w + (len(kpis) - 1) * kpi_gap
    start_x = (SLIDE_W - total) // 2
    for i, (val, lbl) in enumerate(kpis):
        x = start_x + i * (kpi_w + kpi_gap)
        _add_kpi_card(slide, x, kpi_y, kpi_w, kpi_h, val, lbl)

    # Theme cards (2 rows x 3)
    themes = [
        ("Architecture v3", "Spec 051", AI_ACCENT),
        ("IDE-Adapted Mirrors", "Spec 053", SEC_BLUE),
        ("Observability", "Specs 042-045", SEC_PURPLE),
        ("Security Hardening", "Spec 049", AI_ERROR),
        ("Testing Excellence", "Spec 052", AI_SUCCESS),
        ("CI/CD & Release", "Spec 046", AI_WARNING),
    ]
    tc_w = Inches(3.4)
    tc_h = Inches(0.7)
    tc_gap_x = Inches(0.25)
    tc_gap_y = Inches(0.2)
    row_total = 3 * tc_w + 2 * tc_gap_x
    tc_start_x = (SLIDE_W - row_total) // 2
    tc_start_y = Inches(4.2)

    for i, (name, spec, accent) in enumerate(themes):
        col = i % 3
        row = i // 3
        x = tc_start_x + col * (tc_w + tc_gap_x)
        y = tc_start_y + row * (tc_h + tc_gap_y)
        add_card(
            slide,
            x,
            y,
            tc_w,
            tc_h,
            fill_color=AI_CARD_DARK,
            border_color=AI_BORDER_DARK,
            left_accent_color=accent,
        )
        add_textbox(
            slide,
            x + Inches(0.2),
            y + Inches(0.08),
            tc_w - Inches(0.4),
            Inches(0.3),
            text=name,
            font_size=Pt(14),
            color=AI_TEXT_PRIMARY,
            bold=True,
        )
        add_textbox(
            slide,
            x + Inches(0.2),
            y + Inches(0.38),
            tc_w - Inches(0.4),
            Inches(0.25),
            text=spec,
            font_size=Pt(11),
            color=AI_NEUTRAL,
        )

    set_notes(slide, NOTES[2])


# ---------------------------------------------------------------------------
# Slide 3 — Architecture v3
# ---------------------------------------------------------------------------
def slide_03_architecture(prs):
    slide = _blank_slide(prs)
    add_slide_header(slide, "Architecture v3", subtitle="Spec 051 — Clean-sheet redesign")

    # Before → After comparison
    col_w = Inches(5.0)
    col_gap = Inches(0.8)
    total = 2 * col_w + col_gap
    start_x = (SLIDE_W - total) // 2
    col_y = Inches(2.0)

    # Before column
    add_textbox(
        slide,
        start_x,
        col_y,
        col_w,
        Inches(0.4),
        text="BEFORE",
        font_name=FONT_TITLE,
        font_size=Pt(16),
        color=AI_TEXT_MUTED,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    before_items = [
        "7 agents (flat, some missing)",
        "35 skills (5 were stubs)",
        "No proactive governance",
        "No developer onboarding",
        "No SRE/ops agent",
        "Manual self-improvement",
    ]
    for i, item in enumerate(before_items):
        y = col_y + Inches(0.5) + i * Inches(0.5)
        add_card(
            slide,
            start_x,
            y,
            col_w,
            Inches(0.42),
            fill_color=AI_CARD_DARK,
            border_color=AI_BORDER_DARK,
        )
        add_textbox(
            slide,
            start_x + Inches(0.15),
            y + Inches(0.05),
            col_w - Inches(0.3),
            Inches(0.32),
            text=f"  {item}",
            font_size=Pt(12),
            color=AI_TEXT_MUTED,
        )

    # Arrow
    arrow_x = start_x + col_w + Inches(0.1)
    arrow = slide.shapes.add_shape(55, arrow_x, Inches(3.5), Inches(0.6), Inches(0.5))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = AI_ACCENT
    arrow.line.fill.background()

    # After column
    after_x = start_x + col_w + col_gap
    add_textbox(
        slide,
        after_x,
        col_y,
        col_w,
        Inches(0.4),
        text="AFTER",
        font_name=FONT_TITLE,
        font_size=Pt(16),
        color=AI_ACCENT,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    after_items = [
        ("10 agents (3-tier hierarchy)", AI_ACCENT),
        ("40 skills (all functional)", AI_ACCENT),
        ("Guard agent — proactive governance", AI_SUCCESS),
        ("Guide agent — developer growth", SEC_BLUE),
        ("Operate agent — SRE & runbooks", SEC_PURPLE),
        ("Evolve skill — self-improvement loop", AI_WARNING),
    ]
    for i, (item, accent) in enumerate(after_items):
        y = col_y + Inches(0.5) + i * Inches(0.5)
        add_card(
            slide,
            after_x,
            y,
            col_w,
            Inches(0.42),
            fill_color=AI_CARD_DARK,
            border_color=AI_BORDER_DARK,
            left_accent_color=accent,
        )
        add_textbox(
            slide,
            after_x + Inches(0.2),
            y + Inches(0.05),
            col_w - Inches(0.4),
            Inches(0.32),
            text=f"  {item}",
            font_size=Pt(12),
            color=AI_TEXT_LIGHT,
        )

    # Rename callout at bottom
    add_textbox(
        slide,
        LEFT_MARGIN,
        Inches(5.6),
        CONTENT_W,
        Inches(0.35),
        text=(
            "12 skills renamed for clarity: build>code "
            "cicd>pipeline db>schema feature-gap>gap perf>performance"
        ),
        font_size=Pt(11),
        color=AI_NEUTRAL,
        alignment=PP_ALIGN.CENTER,
    )

    set_notes(slide, NOTES[3])


# ---------------------------------------------------------------------------
# Slide 4 — IDE-Adapted Mirrors
# ---------------------------------------------------------------------------
def slide_04_ide_mirrors(prs):
    slide = _blank_slide(prs)
    add_slide_header(slide, "IDE-Adapted Mirrors", subtitle="Spec 053 — Full behavior in every IDE")

    # Problem statement
    add_card(
        slide,
        LEFT_MARGIN,
        Inches(2.0),
        CONTENT_W,
        Inches(1.0),
        fill_color=AI_CARD_DARK,
        border_color=AI_ERROR,
    )
    add_textbox(
        slide,
        LEFT_MARGIN + Inches(0.2),
        Inches(2.1),
        CONTENT_W - Inches(0.4),
        Inches(0.3),
        text="PROBLEM",
        font_name=FONT_TITLE,
        font_size=Pt(14),
        color=AI_ERROR,
        bold=True,
    )
    add_textbox(
        slide,
        LEFT_MARGIN + Inches(0.2),
        Inches(2.45),
        CONTENT_W - Inches(0.4),
        Inches(0.45),
        text=(
            "IDE thin wrappers reference .ai-engineering/ but most IDEs "
            "don't execute references. Behaviors like Interrogation Phase never run."
        ),
        font_size=Pt(13),
        color=AI_TEXT_PRIMARY,
    )

    # Solution
    add_card(
        slide,
        LEFT_MARGIN,
        Inches(3.3),
        CONTENT_W,
        Inches(1.0),
        fill_color=AI_CARD_DARK,
        border_color=AI_SUCCESS,
    )
    add_textbox(
        slide,
        LEFT_MARGIN + Inches(0.2),
        Inches(3.4),
        CONTENT_W - Inches(0.4),
        Inches(0.3),
        text="SOLUTION",
        font_name=FONT_TITLE,
        font_size=Pt(14),
        color=AI_SUCCESS,
        bold=True,
    )
    add_textbox(
        slide,
        LEFT_MARGIN + Inches(0.2),
        Inches(3.75),
        CONTENT_W - Inches(0.4),
        Inches(0.45),
        text=(
            "Full IDE-adapted copies in .claude/, .github/, .agents/ "
            "-- each IDE gets complete, executable content"
        ),
        font_size=Pt(13),
        color=AI_TEXT_PRIMARY,
    )

    # Migration phases
    phase_y = Inches(4.7)
    phase_w = Inches(1.1)
    phase_h = Inches(0.7)
    phase_gap = Inches(0.1)
    phases = [
        ("P0", "Scaffold"),
        ("P1", "Sync\nscript"),
        ("P2", "Delete\nlegacy"),
        ("P3-4", "Install\n+ owners"),
        ("P5", "Update\nvalidators"),
        ("P6-7", "Skills\nservice"),
        ("P8", "Legacy\nmigration"),
        ("P9", "Update\ntests"),
    ]
    total = len(phases) * phase_w + (len(phases) - 1) * phase_gap
    start_x = (SLIDE_W - total) // 2

    for i, (label, desc) in enumerate(phases):
        x = start_x + i * (phase_w + phase_gap)
        fill = AI_SUCCESS if i <= 7 else AI_PRIMARY_LIGHT
        add_rounded_rect(slide, x, phase_y, phase_w, phase_h, fill_color=fill, border_color=None)
        add_textbox(
            slide,
            x,
            phase_y + Inches(0.05),
            phase_w,
            Inches(0.22),
            text=label,
            font_name=FONT_TITLE,
            font_size=Pt(11),
            color=AI_WHITE,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )
        add_textbox(
            slide,
            x,
            phase_y + Inches(0.28),
            phase_w,
            Inches(0.4),
            text=desc,
            font_size=Pt(9),
            color=AI_WHITE,
            alignment=PP_ALIGN.CENTER,
        )

    add_textbox(
        slide,
        LEFT_MARGIN,
        Inches(5.65),
        CONTENT_W,
        Inches(0.3),
        text="Canonical source: src/ai_engineering/templates/.ai-engineering/",
        font_size=Pt(11),
        color=AI_NEUTRAL,
        alignment=PP_ALIGN.CENTER,
    )

    set_notes(slide, NOTES[4])


# ---------------------------------------------------------------------------
# Slide 5 — Observability & Telemetry
# ---------------------------------------------------------------------------
def slide_05_observability(prs):
    slide = _blank_slide(prs)
    add_slide_header(
        slide,
        "Observability & Telemetry",
        subtitle="Specs 042, 044, 045 — Rich dashboards + cross-IDE tracking",
    )

    # 5 dashboard cards
    dash_y = Inches(2.2)
    dash_w = Inches(2.0)
    dash_h = Inches(1.8)
    dash_gap = Inches(0.15)
    dashboards = [
        ("Engineer", "Code quality\nGate performance\nDecision health", AI_ACCENT),
        ("Team", "Token economy\nNoise ratio\nSkill usage", SEC_BLUE),
        ("AI", "Context efficiency\nSkill efficiency\nPattern detection", SEC_PURPLE),
        ("DORA", "Lead time\nChange failure\nDeploy frequency", AI_WARNING),
        ("Health", "Signal quality\nTest confidence\nSecurity posture", AI_SUCCESS),
    ]
    total = len(dashboards) * dash_w + (len(dashboards) - 1) * dash_gap
    start_x = (SLIDE_W - total) // 2

    for i, (name, metrics, accent) in enumerate(dashboards):
        x = start_x + i * (dash_w + dash_gap)
        add_card(
            slide,
            x,
            dash_y,
            dash_w,
            dash_h,
            fill_color=AI_CARD_DARK,
            border_color=AI_BORDER_DARK,
            left_accent_color=accent,
        )
        add_textbox(
            slide,
            x + Inches(0.18),
            dash_y + Inches(0.1),
            dash_w - Inches(0.3),
            Inches(0.3),
            text=name,
            font_name=FONT_TITLE,
            font_size=Pt(14),
            color=accent,
            bold=True,
        )
        add_textbox(
            slide,
            x + Inches(0.18),
            dash_y + Inches(0.5),
            dash_w - Inches(0.3),
            Inches(1.2),
            text=metrics,
            font_size=Pt(11),
            color=AI_TEXT_MUTED,
        )

    # Features row
    feat_y = Inches(4.5)
    feat_w = Inches(3.3)
    feat_h = Inches(1.2)
    feat_gap = Inches(0.25)
    features = [
        (
            "Cross-IDE Telemetry",
            [
                "• ai-eng signals emit skill_invoked",
                "• ai-eng signals emit agent_dispatched",
                "• VCS context: branch, SHA, repo URL",
            ],
        ),
        (
            "Rich Formatting",
            [
                "• Progress bars with percentages",
                "• Score badges (green/yellow/red)",
                "• Metric tables with color coding",
            ],
        ),
        (
            "Dual Output",
            [
                "• Human: Rich-formatted tables",
                "• Machine: --json with SuccessEnvelope",
                "• 4 new cli_ui primitives",
            ],
        ),
    ]
    total_f = len(features) * feat_w + (len(features) - 1) * feat_gap
    start_xf = (SLIDE_W - total_f) // 2

    for i, (title, bullets) in enumerate(features):
        x = start_xf + i * (feat_w + feat_gap)
        _add_feature_card(slide, x, feat_y, feat_w, feat_h, title, bullets)

    set_notes(slide, NOTES[5])


# ---------------------------------------------------------------------------
# Slide 6 — Security Hardening
# ---------------------------------------------------------------------------
def slide_06_security(prs):
    slide = _blank_slide(prs)
    add_slide_header(
        slide, "Security Hardening", subtitle="Spec 049 — Zero suppression, zero tolerance"
    )

    # Vulnerability remediation
    vuln_x = LEFT_MARGIN
    vuln_y = Inches(2.2)
    vuln_w = Inches(5.2)
    vuln_h = Inches(2.3)

    add_card(
        slide, vuln_x, vuln_y, vuln_w, vuln_h, fill_color=AI_CARD_DARK, border_color=AI_BORDER_DARK
    )
    add_textbox(
        slide,
        vuln_x + Inches(0.2),
        vuln_y + Inches(0.1),
        vuln_w - Inches(0.4),
        Inches(0.3),
        text="SonarCloud Remediation",
        font_name=FONT_TITLE,
        font_size=Pt(16),
        color=AI_ERROR,
        bold=True,
    )

    vulns = [
        {
            "text": "  4x BLOCKER — Path traversal (S2083) → path validation",
            "font_size": Pt(12),
            "color": AI_TEXT_PRIMARY,
        },
        {
            "text": "  1x MAJOR — Command injection (S6350) → arg validation",
            "font_size": Pt(12),
            "color": AI_TEXT_PRIMARY,
        },
        {
            "text": "  4x Regex DoS hotspots (S5852) → verified safe",
            "font_size": Pt(12),
            "color": AI_TEXT_MUTED,
        },
        {"text": "", "font_size": Pt(8), "color": AI_TEXT_MUTED},
        {
            "text": "  All fixed by code changes — zero suppression comments",
            "font_size": Pt(13),
            "color": AI_SUCCESS,
            "bold": True,
        },
    ]
    add_rich_textbox(
        slide,
        vuln_x + Inches(0.2),
        vuln_y + Inches(0.55),
        vuln_w - Inches(0.4),
        Inches(1.6),
        vulns,
        line_spacing=Pt(12),
    )

    # Gate pipeline
    gate_x = vuln_x + vuln_w + Inches(0.4)
    gate_y = Inches(2.2)
    gate_w = Inches(5.2)

    add_textbox(
        slide,
        gate_x,
        gate_y,
        gate_w,
        Inches(0.35),
        text="4-Stage Gate Pipeline",
        font_name=FONT_TITLE,
        font_size=Pt(16),
        color=AI_ACCENT,
        bold=True,
    )

    gates = [
        ("Pre-commit", "gitleaks, ruff format/check", AI_ACCENT),
        ("Commit-msg", "Message policy, trailers", SEC_BLUE),
        ("Pre-push", "semgrep, pip-audit, pytest, ty", SEC_PURPLE),
        ("CI", "Full matrix + SonarCloud", AI_WARNING),
    ]
    for i, (name, checks, accent) in enumerate(gates):
        gy = gate_y + Inches(0.5) + i * Inches(0.55)
        add_card(
            slide,
            gate_x,
            gy,
            gate_w,
            Inches(0.48),
            fill_color=AI_CARD_DARK,
            border_color=AI_BORDER_DARK,
            left_accent_color=accent,
        )
        add_textbox(
            slide,
            gate_x + Inches(0.2),
            gy + Inches(0.05),
            Inches(1.5),
            Inches(0.35),
            text=name,
            font_size=Pt(12),
            color=accent,
            bold=True,
        )
        add_textbox(
            slide,
            gate_x + Inches(1.8),
            gy + Inches(0.05),
            gate_w - Inches(2.0),
            Inches(0.35),
            text=checks,
            font_size=Pt(11),
            color=AI_TEXT_MUTED,
        )

    # New rule callout
    add_card(
        slide,
        LEFT_MARGIN,
        Inches(4.9),
        CONTENT_W,
        Inches(0.7),
        fill_color=AI_CARD_DARK,
        border_color=AI_WARNING,
    )
    add_textbox(
        slide,
        LEFT_MARGIN + Inches(0.2),
        Inches(4.98),
        CONTENT_W - Inches(0.4),
        Inches(0.25),
        text="NEW GOVERNANCE RULE",
        font_name=FONT_TITLE,
        font_size=Pt(12),
        color=AI_WARNING,
        bold=True,
    )
    add_textbox(
        slide,
        LEFT_MARGIN + Inches(0.2),
        Inches(5.25),
        CONTENT_W - Inches(0.4),
        Inches(0.25),
        text=(
            "No suppression comments allowed (NOSONAR, nosec, noqa, "
            "type: ignore, pragma: no cover). Fix root cause or escalate."
        ),
        font_size=Pt(12),
        color=AI_TEXT_PRIMARY,
    )

    # Snyk badge
    add_textbox(
        slide,
        LEFT_MARGIN,
        Inches(5.85),
        CONTENT_W,
        Inches(0.3),
        text="+ Snyk CI/CD integration for continuous dependency vulnerability monitoring",
        font_size=Pt(11),
        color=AI_NEUTRAL,
        alignment=PP_ALIGN.CENTER,
    )

    set_notes(slide, NOTES[6])


# ---------------------------------------------------------------------------
# Slide 7 — Testing Excellence
# ---------------------------------------------------------------------------
def slide_07_testing(prs):
    slide = _blank_slide(prs)
    add_slide_header(
        slide, "Testing Excellence", subtitle="Spec 052 — TDD protocol + Interrogation Phase"
    )

    col_w = Inches(3.4)
    col_gap = Inches(0.25)
    total = 3 * col_w + 2 * col_gap
    start_x = (SLIDE_W - total) // 2
    col_y = Inches(2.0)

    # TDD Protocol
    _add_feature_card(
        slide,
        start_x,
        col_y,
        col_w,
        Inches(2.8),
        "TDD Protocol",
        [
            "• RED phase: write failing tests first",
            "• GREEN phase: implement to pass",
            "• Tests are immutable during impl",
            "• Iron Law: never weaken tests",
            "• Writer/implementer separation",
            "• Enforced by dispatch protocol",
        ],
        accent=AI_ERROR,
    )

    # Interrogation Phase
    _add_feature_card(
        slide,
        start_x + col_w + col_gap,
        col_y,
        col_w,
        Inches(2.8),
        "Interrogation Phase",
        [
            "• Challenge assumptions before planning",
            "• Map: KNOWN / ASSUMED / UNKNOWN",
            "• Executable acceptance criteria",
            "• Verification Commands with expected output",
            "• Replaces prose-based verification",
            "• Built into plan agent workflow",
        ],
        accent=SEC_BLUE,
    )

    # Test Infrastructure
    _add_feature_card(
        slide,
        start_x + 2 * (col_w + col_gap),
        col_y,
        col_w,
        Inches(2.8),
        "Test Infrastructure",
        [
            "• 91% coverage (target: 80%)",
            "• 20 stacks supported",
            "• Fakes over mocks pattern",
            "• AAA (Arrange-Act-Assert)",
            "• 63 patches → 1 fixture",
            "• Flaky test diagnostics (6 categories)",
        ],
        accent=AI_SUCCESS,
    )

    # Bottom callout
    add_textbox(
        slide,
        LEFT_MARGIN,
        Inches(5.2),
        CONTENT_W,
        Inches(0.3),
        text=(
            "Test coverage went from 26% (meaningful tests) to 91% this sprint "
            "-- covering agents, skills, and governance logic"
        ),
        font_size=Pt(12),
        color=AI_TEXT_MUTED,
        alignment=PP_ALIGN.CENTER,
    )

    set_notes(slide, NOTES[7])


# ---------------------------------------------------------------------------
# Slide 8 — CI/CD & Release
# ---------------------------------------------------------------------------
def slide_08_cicd(prs):
    slide = _blank_slide(prs)
    add_slide_header(slide, "CI/CD & Release", subtitle="Spec 046 — Zero-rebuild release model")

    # Main feature: Zero-rebuild
    feat_x = LEFT_MARGIN
    feat_y = Inches(2.2)
    feat_w = Inches(5.2)
    feat_h = Inches(2.5)

    add_card(slide, feat_x, feat_y, feat_w, feat_h, fill_color=AI_CARD_DARK, border_color=AI_ACCENT)
    add_textbox(
        slide,
        feat_x + Inches(0.2),
        feat_y + Inches(0.1),
        feat_w - Inches(0.4),
        Inches(0.35),
        text="Zero-Rebuild Release",
        font_name=FONT_TITLE,
        font_size=Pt(18),
        color=AI_ACCENT,
        bold=True,
    )

    flow_lines = [
        {
            "text": "CI builds & validates → artifact stored",
            "font_size": Pt(13),
            "color": AI_TEXT_PRIMARY,
        },
        {"text": "                    ↓", "font_size": Pt(13), "color": AI_ACCENT},
        {
            "text": "Release downloads exact CI artifact",
            "font_size": Pt(13),
            "color": AI_TEXT_PRIMARY,
        },
        {"text": "                    ↓", "font_size": Pt(13), "color": AI_ACCENT},
        {"text": "PyPI + GitHub Releases publish", "font_size": Pt(13), "color": AI_TEXT_PRIMARY},
        {"text": "", "font_size": Pt(6), "color": AI_TEXT_MUTED},
        {
            "text": "What CI tests = what ships. Zero drift.",
            "font_size": Pt(14),
            "color": AI_SUCCESS,
            "bold": True,
        },
    ]
    add_rich_textbox(
        slide,
        feat_x + Inches(0.2),
        feat_y + Inches(0.6),
        feat_w - Inches(0.4),
        Inches(1.8),
        flow_lines,
        line_spacing=Pt(8),
    )

    # Other improvements
    other_x = feat_x + feat_w + Inches(0.4)
    other_y = Inches(2.2)
    other_w = Inches(5.2)

    improvements = [
        (
            "Dependabot Grouping",
            "PRs grouped, exempted from trailer\nchecks to reduce noise",
            AI_ACCENT,
        ),
        ("CI Matrix", "Linux + macOS + Windows\nPython 3.11+ · Stabilizing", SEC_BLUE),
        (
            "8 Automated Workflows",
            "CI fixer, code simplifier, daily triage,\nPR review, security scan, weekly health",
            SEC_PURPLE,
        ),
        ("Install Smoke Tests", "Framework installation validation\nacross platforms", AI_WARNING),
    ]
    card_h = Inches(0.9)
    card_gap = Inches(0.15)
    for i, (title, desc, accent) in enumerate(improvements):
        y = other_y + i * (card_h + card_gap)
        add_card(
            slide,
            other_x,
            y,
            other_w,
            card_h,
            fill_color=AI_CARD_DARK,
            border_color=AI_BORDER_DARK,
            left_accent_color=accent,
        )
        add_textbox(
            slide,
            other_x + Inches(0.2),
            y + Inches(0.08),
            other_w - Inches(0.4),
            Inches(0.25),
            text=title,
            font_size=Pt(13),
            color=accent,
            bold=True,
        )
        add_textbox(
            slide,
            other_x + Inches(0.2),
            y + Inches(0.4),
            other_w - Inches(0.4),
            Inches(0.45),
            text=desc,
            font_size=Pt(11),
            color=AI_TEXT_MUTED,
        )

    set_notes(slide, NOTES[8])


# ---------------------------------------------------------------------------
# Slide 9 — Quality Metrics
# ---------------------------------------------------------------------------
def slide_09_quality(prs):
    slide = _blank_slide(prs)
    add_slide_header(slide, "Quality Metrics", subtitle="Enforced — not aspirational")

    # Quality metrics table
    tbl, _shape = add_styled_table(
        slide,
        LEFT_MARGIN,
        Inches(2.2),
        CONTENT_W,
        Inches(3.5),
        7,
        4,
        header_fill=AI_ACCENT,
    )

    headers = ["Metric", "Target", "Current", "Status"]
    data = [
        ("Test Coverage", "≥ 80%", "91%", "EXCEEDED"),
        ("Cyclomatic Complexity", "≤ 10/function", "Enforced (ruff)", "PASSING"),
        ("Cognitive Complexity", "≤ 15/function", "Enforced (ruff)", "PASSING"),
        ("Gate Pass Rate", "100%", "100%", "PASSING"),
        ("Token Efficiency", "≥ 95% deferred", "99.19%", "EXCEEDED"),
        ("Code Duplication", "≤ 3%", "Measured", "TRACKING"),
    ]

    for ci, h in enumerate(headers):
        _style_cell(
            tbl.cell(0, ci),
            h,
            bold=True,
            color=AI_BG_DARK,
            font_size=Pt(13),
            alignment=PP_ALIGN.CENTER,
        )

    for ri, (metric, target, current, status) in enumerate(data, start=1):
        bg = AI_CARD_DARK if ri % 2 == 0 else AI_BG_DARK
        _style_cell(tbl.cell(ri, 0), metric, font_size=Pt(12), fill_color=bg, color=AI_TEXT_PRIMARY)
        _style_cell(
            tbl.cell(ri, 1),
            target,
            font_size=Pt(12),
            fill_color=bg,
            color=AI_TEXT_MUTED,
            alignment=PP_ALIGN.CENTER,
        )
        _style_cell(
            tbl.cell(ri, 2),
            current,
            font_size=Pt(12),
            fill_color=bg,
            color=AI_WHITE,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )
        status_color = AI_SUCCESS if status in ("EXCEEDED", "PASSING") else AI_WARNING
        _style_cell(
            tbl.cell(ri, 3),
            status,
            font_size=Pt(12),
            fill_color=bg,
            color=status_color,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )

    # Bottom note
    add_textbox(
        slide,
        LEFT_MARGIN,
        Inches(6.0),
        CONTENT_W,
        Inches(0.3),
        text="All metrics enforced by 4-stage gate pipeline — violations block the commit/push",
        font_size=Pt(12),
        color=AI_NEUTRAL,
        alignment=PP_ALIGN.CENTER,
    )

    set_notes(slide, NOTES[9])


# ---------------------------------------------------------------------------
# Slide 10 — Governance Surface
# ---------------------------------------------------------------------------
def slide_10_governance(prs):
    slide = _blank_slide(prs)
    add_slide_header(slide, "Governance Surface", subtitle="38 specs delivered — 4 IDEs supported")

    # Big numbers row
    kpi_y = Inches(2.2)
    kpi_w = Inches(2.0)
    kpi_h = Inches(1.3)
    kpi_gap = Inches(0.25)
    kpis = [
        ("10", "Agents"),
        ("40", "Skills"),
        ("50+", "Standards"),
        ("38", "Specs Delivered"),
        ("4", "IDEs"),
    ]
    total = len(kpis) * kpi_w + (len(kpis) - 1) * kpi_gap
    start_x = (SLIDE_W - total) // 2

    for i, (val, lbl) in enumerate(kpis):
        x = start_x + i * (kpi_w + kpi_gap)
        _add_kpi_card(slide, x, kpi_y, kpi_w, kpi_h, val, lbl)

    # IDE support cards
    ide_y = Inches(4.0)
    ide_w = Inches(2.5)
    ide_h = Inches(1.2)
    ide_gap = Inches(0.2)
    ides = [
        ("Claude Code", "CLAUDE.md + 40 skills\n+ 8 agents", AI_ACCENT),
        ("GitHub Copilot", "copilot-instructions.md\n+ prompt files + agents", SEC_BLUE),
        ("Gemini CLI", "GEMINI.md\n+ adapted instructions", SEC_PURPLE),
        ("OpenAI Codex", "AGENTS.md\n+ native format", AI_WARNING),
    ]
    total_i = len(ides) * ide_w + (len(ides) - 1) * ide_gap
    start_xi = (SLIDE_W - total_i) // 2

    for i, (name, desc, accent) in enumerate(ides):
        x = start_xi + i * (ide_w + ide_gap)
        add_card(
            slide,
            x,
            ide_y,
            ide_w,
            ide_h,
            fill_color=AI_CARD_DARK,
            border_color=AI_BORDER_DARK,
            left_accent_color=accent,
        )
        add_textbox(
            slide,
            x + Inches(0.18),
            ide_y + Inches(0.1),
            ide_w - Inches(0.3),
            Inches(0.3),
            text=name,
            font_name=FONT_TITLE,
            font_size=Pt(13),
            color=accent,
            bold=True,
        )
        add_textbox(
            slide,
            x + Inches(0.18),
            ide_y + Inches(0.5),
            ide_w - Inches(0.3),
            Inches(0.65),
            text=desc,
            font_size=Pt(11),
            color=AI_TEXT_MUTED,
        )

    # State files
    add_textbox(
        slide,
        LEFT_MARGIN,
        Inches(5.5),
        CONTENT_W,
        Inches(0.3),
        text=(
            "State: 17 decisions (SHA-256 hashed)  |  185+ audit events  "
            "|  Ownership map  |  Sources lock"
        ),
        font_size=Pt(12),
        color=AI_NEUTRAL,
        alignment=PP_ALIGN.CENTER,
    )

    set_notes(slide, NOTES[10])


# ---------------------------------------------------------------------------
# Slide 11 — Risks & Next Sprint
# ---------------------------------------------------------------------------
def slide_11_risks_next(prs):
    slide = _blank_slide(prs)
    add_slide_header(slide, "Risks & Next Sprint")

    # Risks section
    risk_x = LEFT_MARGIN
    risk_y = Inches(2.0)
    risk_w = Inches(5.2)

    add_textbox(
        slide,
        risk_x,
        risk_y,
        risk_w,
        Inches(0.35),
        text="ACTIVE RISKS",
        font_name=FONT_TITLE,
        font_size=Pt(16),
        color=AI_ERROR,
        bold=True,
    )

    risks = [
        (
            "MEDIUM",
            "ty type checker: 84 diagnostics",
            "Non-blocking tech debt — systematic remediation planned",
        ),
        (
            "MEDIUM",
            "Cross-OS CI stabilization",
            "Windows + macOS need further hardening in CI matrix",
        ),
        (
            "LOW",
            "Architecture v3 agent testing",
            "3 new agents (guard, guide, operate) need integration tests",
        ),
    ]
    for i, (severity, title, desc) in enumerate(risks):
        ry = risk_y + Inches(0.5) + i * Inches(0.85)
        sev_color = AI_WARNING if severity == "MEDIUM" else AI_ACCENT
        add_card(
            slide,
            risk_x,
            ry,
            risk_w,
            Inches(0.75),
            fill_color=AI_CARD_DARK,
            border_color=AI_BORDER_DARK,
            left_accent_color=sev_color,
        )
        add_textbox(
            slide,
            risk_x + Inches(0.2),
            ry + Inches(0.08),
            Inches(1.0),
            Inches(0.25),
            text=severity,
            font_name=FONT_TITLE,
            font_size=Pt(10),
            color=sev_color,
            bold=True,
        )
        add_textbox(
            slide,
            risk_x + Inches(1.2),
            ry + Inches(0.08),
            risk_w - Inches(1.4),
            Inches(0.25),
            text=title,
            font_size=Pt(12),
            color=AI_TEXT_PRIMARY,
            bold=True,
        )
        add_textbox(
            slide,
            risk_x + Inches(0.2),
            ry + Inches(0.38),
            risk_w - Inches(0.4),
            Inches(0.3),
            text=desc,
            font_size=Pt(10),
            color=AI_TEXT_MUTED,
        )

    # Next Sprint section
    next_x = risk_x + risk_w + Inches(0.4)
    next_y = Inches(2.0)
    next_w = Inches(5.2)

    add_textbox(
        slide,
        next_x,
        next_y,
        next_w,
        Inches(0.35),
        text="NEXT SPRINT",
        font_name=FONT_TITLE,
        font_size=Pt(16),
        color=AI_SUCCESS,
        bold=True,
    )

    nexts = [
        ("Spec 053 Completion", "Finalize Phase 10+ of IDE-adapted mirrors migration", AI_ACCENT),
        (
            "Release Readiness",
            "Docs site (Nextra), PyPI stable release, remote skill sources",
            SEC_BLUE,
        ),
        (
            "Framework v2 Restructure",
            "Complete 8-dimension platform audit remediation (Spec 050)",
            SEC_PURPLE,
        ),
        (
            "Multi-IDE Verification",
            "End-to-end validation across Claude, Copilot, Gemini, Codex",
            AI_WARNING,
        ),
    ]
    for i, (title, desc, accent) in enumerate(nexts):
        ny = next_y + Inches(0.5) + i * Inches(0.85)
        add_card(
            slide,
            next_x,
            ny,
            next_w,
            Inches(0.75),
            fill_color=AI_CARD_DARK,
            border_color=AI_BORDER_DARK,
            left_accent_color=accent,
        )
        add_textbox(
            slide,
            next_x + Inches(0.2),
            ny + Inches(0.08),
            next_w - Inches(0.4),
            Inches(0.25),
            text=title,
            font_size=Pt(13),
            color=accent,
            bold=True,
        )
        add_textbox(
            slide,
            next_x + Inches(0.2),
            ny + Inches(0.38),
            next_w - Inches(0.4),
            Inches(0.3),
            text=desc,
            font_size=Pt(11),
            color=AI_TEXT_MUTED,
        )

    set_notes(slide, NOTES[11])


# ---------------------------------------------------------------------------
# Slide 12 — Q&A
# ---------------------------------------------------------------------------
def slide_12_qa(prs):
    slide = _blank_slide(prs)

    add_accent_bar(slide, Inches(0), Inches(0.4), SLIDE_W, height=Pt(4))

    add_textbox(
        slide,
        LEFT_MARGIN,
        Inches(2.5),
        CONTENT_W,
        Inches(1.0),
        text="Thank You",
        font_name=FONT_TITLE,
        font_size=Pt(52),
        color=AI_WHITE,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )

    bar_w = Inches(3)
    bar_left = (SLIDE_W - bar_w) // 2
    add_accent_bar(slide, bar_left, Inches(3.5), bar_w, height=Pt(3))

    add_textbox(
        slide,
        LEFT_MARGIN,
        Inches(3.8),
        CONTENT_W,
        Inches(0.5),
        text="Questions?",
        font_size=Pt(24),
        color=AI_ACCENT,
        alignment=PP_ALIGN.CENTER,
    )

    # Summary badges
    badge_y = Inches(5.0)
    badge_w = Inches(2.2)
    badge_h = Inches(0.5)
    gap = Inches(0.3)
    badges = [
        "Architecture v3",
        "IDE Mirrors",
        "Observability",
        "Security",
    ]
    total_w = len(badges) * badge_w + (len(badges) - 1) * gap
    start_x = (SLIDE_W - total_w) // 2
    for i, label in enumerate(badges):
        x = start_x + i * (badge_w + gap)
        add_rounded_rect(
            slide, x, badge_y, badge_w, badge_h, fill_color=AI_CARD_DARK, border_color=AI_ACCENT
        )
        add_textbox(
            slide,
            x,
            badge_y,
            badge_w,
            badge_h,
            text=label,
            font_size=Pt(12),
            color=AI_ACCENT,
            bold=True,
            alignment=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )

    add_textbox(
        slide,
        LEFT_MARGIN,
        Inches(5.8),
        CONTENT_W,
        Inches(0.4),
        text="Sprint Feb 16 — Mar 16, 2026  ·  ai-engineering",
        font_size=Pt(14),
        color=AI_NEUTRAL,
        alignment=PP_ALIGN.CENTER,
    )

    add_accent_bar(slide, Inches(0), Inches(7.0), SLIDE_W, height=Pt(4))
    set_notes(slide, NOTES[12])


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_01_title(prs)
    slide_02_overview(prs)
    slide_03_architecture(prs)
    slide_04_ide_mirrors(prs)
    slide_05_observability(prs)
    slide_06_security(prs)
    slide_07_testing(prs)
    slide_08_cicd(prs)
    slide_09_quality(prs)
    slide_10_governance(prs)
    slide_11_risks_next(prs)
    slide_12_qa(prs)

    out = Path(__file__).parent / "sprint-review-2026-03.pptx"
    prs.save(str(out))
    print(f"Generated: {out}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
