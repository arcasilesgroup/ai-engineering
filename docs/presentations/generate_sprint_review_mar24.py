"""Generate Sprint Review presentation (.pptx) - Mar 16 - Mar 24, 2026.

Output: docs/presentations/sprint-review-2026-03-24.pptx
Brand: ai-engineering dark mode — dark backgrounds, teal accents, technical tone.
Audience: engineers, heads, managers — detailed per feature, non-technical language.
"""

from __future__ import annotations

import contextlib
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# ai-engineering Dark-Mode Colour Palette
# ---------------------------------------------------------------------------
AI_BG_DARK = RGBColor(0x0B, 0x11, 0x20)
AI_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
AI_TEXT_PRIMARY = RGBColor(0xE2, 0xE8, 0xF0)

AI_ACCENT = RGBColor(0x00, 0xD4, 0xAA)
AI_PRIMARY = RGBColor(0x1E, 0x3A, 0x5F)
AI_PRIMARY_LIGHT = RGBColor(0x2A, 0x4F, 0x7A)
AI_ERROR = RGBColor(0xEF, 0x44, 0x44)
AI_SUCCESS = RGBColor(0x10, 0xB9, 0x81)
AI_WARNING = RGBColor(0xF5, 0x9E, 0x0B)

AI_TEXT_LIGHT = RGBColor(0xF8, 0xFA, 0xFB)
AI_NEUTRAL = RGBColor(0x64, 0x74, 0x8B)
AI_TEXT_MUTED = RGBColor(0x94, 0xA3, 0xB8)

AI_BORDER_DARK = RGBColor(0x1A, 0x2A, 0x40)
AI_CARD_DARK = RGBColor(0x1E, 0x29, 0x3B)

SEC_BLUE = RGBColor(0x2E, 0x6B, 0xA4)
SEC_BLUE_LIGHT = RGBColor(0x7A, 0xB5, 0xD6)
SEC_TEAL = RGBColor(0x00, 0xD4, 0xAA)
SEC_PURPLE = RGBColor(0x7B, 0x3F, 0xA0)
SEC_PURPLE_LIGHT = RGBColor(0xB3, 0x8B, 0xCF)

# ---------------------------------------------------------------------------
# Typography
# ---------------------------------------------------------------------------
FONT_TITLE = "JetBrains Mono"
FONT_BODY = "Inter"

# ---------------------------------------------------------------------------
# Slide dimensions (16:9)
# ---------------------------------------------------------------------------
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

LEFT_MARGIN = Inches(1.2)
CONTENT_W = SLIDE_W - LEFT_MARGIN - Inches(1.2)

_PT_3 = Pt(3)
_PT_14 = Pt(14)
_PT_20 = Pt(20)
_PT_1 = Pt(1)
_PT_12 = Pt(12)

# ---------------------------------------------------------------------------
# Speaker Notes
# ---------------------------------------------------------------------------
NOTES: dict[int, str] = {
    1: (
        "Bienvenidos a la sprint review del periodo 16 al 24 de marzo de 2026.\n\n"
        "Ha sido una semana de transformacion profunda — 30 commits en 31 PRs "
        "mergeados en 8 dias activos, tocando mas de 2,700 archivos.\n\n"
        "El tema central: simplificacion radical y madurez multi-IDE."
    ),
    2: (
        "Vista general del sprint.\n\n"
        "30 commits, 31 PRs mergeados en 8 dias de trabajo. "
        "151K lineas insertadas, 96K eliminadas — un neto de ~55K lineas. "
        "Pero lo importante es que eliminamos mas de lo que anadimos en el spec-055, "
        "que redujo significativamente la complejidad del framework.\n\n"
        "7 specs completados, casi 2,000 tests, y el framework ahora soporta "
        "37 skills en 4 IDEs."
    ),
    3: (
        "La entrega mas grande del sprint: Radical Simplification, spec-055.\n\n"
        "1,016 archivos cambiados. 74,000 lineas eliminadas contra 53,000 anadidas. "
        "Eliminamos el directorio .ai-engineering/skills/ y .ai-engineering/agents/ completos. "
        "Los IDE mirrors ya SON la unica fuente de skills y agents.\n\n"
        "Se eliminaron 2,175 lineas de codigo de generacion de pipelines programatico. "
        "Ahora /ai-pipeline lee el manifest y genera pipelines con IA, sin codigo custom.\n\n"
        "El resultado: un framework mas ligero, mas facil de mantener, y mas claro."
    ),
    4: (
        "Multi-IDE Parity es el segundo gran tema.\n\n"
        "spec-053 completo los IDE-Adapted Mirrors — contenido completo en .claude/, "
        ".github/, y .agents/ en vez de thin wrappers.\n\n"
        "spec-062 anadio sync de handlers, registro de skills, y separadores "
        "en los archivos .prompt.md concatenados. 125 tests de regresion nuevos.\n\n"
        "Ademas, implementamos hooks de telemetria completos para GitHub Copilot "
        "con paridad con Claude Code — detect de agent names y skill invocations."
    ),
    5: (
        "Pipeline Skill v2 fue una reescritura mayor.\n\n"
        "spec-054 reescribio el generador de CI, anadio advanced patterns para "
        "GitHub Actions y Azure Pipelines, y implemento telemetria via hooks.\n\n"
        "Luego spec-059 dio el paso mas radical: elimino completamente el generador "
        "programatico de pipelines. -2,175 lineas. Ahora /ai-pipeline lee "
        "cicd.standards_url del manifest y genera pipelines con IA.\n\n"
        "Mas simple, mas flexible, y el equipo ya no mantiene templates de pipeline."
    ),
    6: (
        "Enterprise features y mejoras de delivery.\n\n"
        "spec-056 reorganizo el manifest.yml, anadio soporte para artifact feeds "
        "empresariales, y limpio la configuracion.\n\n"
        "spec-058 agrego un watch & fix loop para /ai-pr, "
        "permitiendo re-intentar automaticamente el paso 14 del PR workflow "
        "cuando falla por issues transitorios.\n\n"
        "El skill frontmatter evoluciono con campos de effort level, "
        "permitiendo al framework asignar compute segun la complejidad cognitiva del skill."
    ),
    7: (
        "CI Reliability fue clave para mantener la confianza en el pipeline.\n\n"
        "spec-060 elimino falsos positivos en CI y en los smoke tests de instalacion. "
        "Corrigió el mapeo de tests, fallbacks de test_scope, "
        "y la gestion de Dependabot PRs.\n\n"
        "Varios PRs de calidad: eliminacion de comentarios de supresion, "
        "restauracion de validacion a PASS, actualizacion de documentacion, "
        "y reescritura de test_doctor de 63 patches a 1 fixture."
    ),
    8: (
        "Metricas de calidad del sprint.\n\n"
        "1,968 tests en el repositorio. Gate pass rate 100%%. "
        "37 skills, 9 agents, 4 IDEs soportados.\n\n"
        "El dato mas significativo: -21K lineas netas eliminadas "
        "en spec-055. El framework es mas pequeno y potente que hace una semana."
    ),
    9: (
        "Riesgos activos y prioridades del proximo sprint.\n\n"
        "Los archivos modificados sin commit del git status actual son las nuevas "
        "skills y agents en desarrollo: autopilot, eval, media, slides, video-editing.\n\n"
        "Proximo sprint: completar esas nuevas skills, "
        "validacion cross-IDE end-to-end, y docs site."
    ),
    10: (
        "Gracias. Preguntas?\n\n"
        "Resumen: simplificacion radical (-21K LOC netas), "
        "multi-IDE con paridad completa, pipeline generation via AI, "
        "y CI reliability mejorada."
    ),
}


# ---------------------------------------------------------------------------
# Helper utilities
# ---------------------------------------------------------------------------
def _font(
    run,
    *,
    name: str = FONT_BODY,
    size: Pt = _PT_14,
    color: RGBColor = AI_TEXT_LIGHT,
    bold: bool = False,
):
    run.font.name = name
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    font_name: str = FONT_BODY,
    font_size: Pt = _PT_14,
    color: RGBColor = AI_TEXT_LIGHT,
    bold: bool = False,
    alignment: PP_ALIGN = PP_ALIGN.LEFT,
    word_wrap: bool = True,
    anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = word_wrap
    tf.auto_size = None
    with contextlib.suppress(Exception):
        tf.paragraphs[0].alignment = alignment
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    _font(run, name=font_name, size=font_size, color=color, bold=bold)
    txbox.text_frame.paragraphs[0].space_before = Pt(0)
    txbox.text_frame.paragraphs[0].space_after = Pt(0)
    with contextlib.suppress(Exception):
        txbox.text_frame._txBody.bodyPr.set(
            "anchor",
            {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}.get(
                anchor, "t"
            ),
        )
    return txbox


def add_rich_textbox(
    slide, left, top, width, height, lines, *, alignment=PP_ALIGN.LEFT, line_spacing=_PT_20
):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = alignment
        if line_spacing:
            p.space_after = line_spacing
        run = p.add_run()
        run.text = ln.get("text", "")
        _font(
            run,
            name=ln.get("font_name", FONT_BODY),
            size=ln.get("font_size", Pt(14)),
            color=ln.get("color", AI_TEXT_LIGHT),
            bold=ln.get("bold", False),
        )
    return txbox


def add_accent_bar(slide, left, top, width, *, height=_PT_3, color=AI_ACCENT):
    bar = slide.shapes.add_shape(1, left, top, width, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar


def add_card(
    slide,
    left,
    top,
    width,
    height,
    *,
    fill_color=AI_CARD_DARK,
    border_color=None,
    border_width=_PT_1,
    left_accent_color=None,
    left_accent_width=None,
):
    card = slide.shapes.add_shape(1, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = fill_color
    if border_color:
        card.line.color.rgb = border_color
        card.line.width = border_width
    else:
        card.line.fill.background()
    if left_accent_color:
        _accent_w = left_accent_width if left_accent_width is not None else Inches(0.08)
        accent = slide.shapes.add_shape(1, left, top, _accent_w, height)
        accent.fill.solid()
        accent.fill.fore_color.rgb = left_accent_color
        accent.line.fill.background()
    return card


def add_rounded_rect(
    slide, left, top, width, height, *, fill_color=AI_CARD_DARK, border_color=AI_TEXT_MUTED
):
    rr = slide.shapes.add_shape(5, left, top, width, height)
    rr.fill.solid()
    rr.fill.fore_color.rgb = fill_color
    if border_color:
        rr.line.color.rgb = border_color
        rr.line.width = Pt(1)
    else:
        rr.line.fill.background()
    return rr


def set_notes(slide, text: str):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


def add_slide_header(slide, title: str, *, subtitle: str | None = None):
    add_accent_bar(slide, LEFT_MARGIN, Inches(0.6), CONTENT_W)
    add_textbox(
        slide,
        LEFT_MARGIN,
        Inches(0.75),
        CONTENT_W,
        Inches(0.7),
        text=title,
        font_name=FONT_TITLE,
        font_size=Pt(34),
        color=AI_WHITE,
        bold=True,
    )
    if subtitle:
        add_textbox(
            slide,
            LEFT_MARGIN,
            Inches(1.35),
            CONTENT_W,
            Inches(0.4),
            text=subtitle,
            font_name=FONT_BODY,
            font_size=Pt(18),
            color=AI_NEUTRAL,
        )


def _blank_slide(prs):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = AI_BG_DARK
    return slide


def add_styled_table(slide, left, top, width, height, rows, cols, *, header_fill=AI_ACCENT):
    shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    tbl = shape.table
    for ci in range(cols):
        cell = tbl.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
    return tbl, shape


def _style_cell(
    cell,
    text,
    *,
    font_size=_PT_12,
    bold=False,
    color=AI_TEXT_LIGHT,
    alignment=PP_ALIGN.LEFT,
    font_name=FONT_BODY,
    fill_color=None,
):
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    _font(run, name=font_name, size=font_size, color=color, bold=bold)
    cell.text_frame.word_wrap = True
    if fill_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill_color
    cell.margin_left = Inches(0.08)
    cell.margin_right = Inches(0.08)
    cell.margin_top = Inches(0.04)
    cell.margin_bottom = Inches(0.04)


# ---------------------------------------------------------------------------
# KPI card helper
# ---------------------------------------------------------------------------
def _add_kpi_card(slide, x, y, w, h, value, label, *, accent=AI_ACCENT):
    add_card(slide, x, y, w, h, fill_color=AI_CARD_DARK, border_color=AI_BORDER_DARK)
    add_textbox(
        slide,
        x,
        y + Inches(0.15),
        w,
        Inches(0.6),
        text=value,
        font_name=FONT_TITLE,
        font_size=Pt(36),
        color=accent,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        x,
        y + Inches(0.8),
        w,
        Inches(0.4),
        text=label,
        font_size=Pt(13),
        color=AI_NEUTRAL,
        alignment=PP_ALIGN.CENTER,
    )


# ---------------------------------------------------------------------------
# Feature card helper
# ---------------------------------------------------------------------------
def _add_feature_card(slide, x, y, w, h, title, bullets, *, accent=AI_ACCENT):
    add_card(
        slide,
        x,
        y,
        w,
        h,
        fill_color=AI_CARD_DARK,
        border_color=AI_BORDER_DARK,
        left_accent_color=accent,
    )
    add_textbox(
        slide,
        x + Inches(0.2),
        y + Inches(0.1),
        w - Inches(0.4),
        Inches(0.35),
        text=title,
        font_size=Pt(14),
        color=AI_TEXT_PRIMARY,
        bold=True,
    )
    lines = [{"text": f"  {b}", "font_size": Pt(11), "color": AI_TEXT_MUTED} for b in bullets]
    add_rich_textbox(
        slide,
        x + Inches(0.2),
        y + Inches(0.45),
        w - Inches(0.4),
        h - Inches(0.55),
        lines,
        line_spacing=Pt(6),
    )


# ---------------------------------------------------------------------------
# Slide 1 — Title
# ---------------------------------------------------------------------------
def slide_01_title(prs):
    slide = _blank_slide(prs)

    add_accent_bar(slide, Inches(0), Inches(0.4), SLIDE_W, height=Pt(4))

    add_textbox(
        slide,
        LEFT_MARGIN,
        Inches(1.8),
        CONTENT_W,
        Inches(1.0),
        text="Sprint Review",
        font_name=FONT_TITLE,
        font_size=Pt(52),
        color=AI_WHITE,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        LEFT_MARGIN,
        Inches(2.85),
        CONTENT_W,
        Inches(0.6),
        text="ai-engineering",
        font_name=FONT_TITLE,
        font_size=Pt(28),
        color=AI_ACCENT,
        alignment=PP_ALIGN.CENTER,
    )

    bar_w = Inches(3)
    bar_left = (SLIDE_W - bar_w) // 2
    add_accent_bar(slide, bar_left, Inches(3.6), bar_w, height=Pt(3))

    add_textbox(
        slide,
        LEFT_MARGIN,
        Inches(3.9),
        CONTENT_W,
        Inches(0.5),
        text="March 16 — March 24, 2026",
        font_size=Pt(20),
        color=AI_TEXT_PRIMARY,
        alignment=PP_ALIGN.CENTER,
    )

    # Sprint badges
    badge_y = Inches(5.0)
    badge_w = Inches(2.0)
    badge_h = Inches(0.45)
    gap = Inches(0.3)
    badges = ["30 Commits", "31 PRs", "7 Specs", "~2K Tests"]
    total_w = len(badges) * badge_w + (len(badges) - 1) * gap
    start_x = (SLIDE_W - total_w) // 2
    for i, label in enumerate(badges):
        x = start_x + i * (badge_w + gap)
        add_rounded_rect(
            slide, x, badge_y, badge_w, badge_h, fill_color=AI_CARD_DARK, border_color=AI_ACCENT
        )
        add_textbox(
            slide,
            x,
            badge_y,
            badge_w,
            badge_h,
            text=label,
            font_size=Pt(12),
            color=AI_ACCENT,
            bold=True,
            alignment=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )

    add_accent_bar(slide, Inches(0), Inches(7.0), SLIDE_W, height=Pt(4))
    set_notes(slide, NOTES[1])


# ---------------------------------------------------------------------------
# Slide 2 — Sprint Overview
# ---------------------------------------------------------------------------
def slide_02_overview(prs):
    slide = _blank_slide(prs)
    add_slide_header(slide, "Sprint Overview", subtitle="Mar 16 - Mar 24, 2026  |  8 active days")

    # KPI cards row
    kpi_y = Inches(2.2)
    kpi_w = Inches(2.4)
    kpi_h = Inches(1.3)
    kpi_gap = Inches(0.3)
    kpis = [
        ("30", "Commits"),
        ("55K+", "Net Lines Changed"),
        ("7", "Specs Completed"),
        ("~2K", "Tests"),
    ]
    total = len(kpis) * kpi_w + (len(kpis) - 1) * kpi_gap
    start_x = (SLIDE_W - total) // 2
    for i, (val, lbl) in enumerate(kpis):
        x = start_x + i * (kpi_w + kpi_gap)
        _add_kpi_card(slide, x, kpi_y, kpi_w, kpi_h, val, lbl)

    # Theme cards (2 rows x 3)
    themes = [
        ("Radical Simplification", "Spec 055 — -21K LOC net", AI_ACCENT),
        ("Multi-IDE Parity", "Specs 053, 062 + Copilot Hooks", SEC_BLUE),
        ("Pipeline v2 & AI-Driven", "Specs 054, 059", SEC_PURPLE),
        ("Enterprise Features", "Specs 056, 058", AI_WARNING),
        ("CI Reliability", "Spec 060 + fixes", AI_SUCCESS),
        ("Skill Framework", "Effort levels + frontmatter", AI_ERROR),
    ]
    tc_w = Inches(3.4)
    tc_h = Inches(0.7)
    tc_gap_x = Inches(0.25)
    tc_gap_y = Inches(0.2)
    row_total = 3 * tc_w + 2 * tc_gap_x
    tc_start_x = (SLIDE_W - row_total) // 2
    tc_start_y = Inches(4.2)

    for i, (name, spec, accent) in enumerate(themes):
        col = i % 3
        row = i // 3
        x = tc_start_x + col * (tc_w + tc_gap_x)
        y = tc_start_y + row * (tc_h + tc_gap_y)
        add_card(
            slide,
            x,
            y,
            tc_w,
            tc_h,
            fill_color=AI_CARD_DARK,
            border_color=AI_BORDER_DARK,
            left_accent_color=accent,
        )
        add_textbox(
            slide,
            x + Inches(0.2),
            y + Inches(0.08),
            tc_w - Inches(0.4),
            Inches(0.3),
            text=name,
            font_size=Pt(14),
            color=AI_TEXT_PRIMARY,
            bold=True,
        )
        add_textbox(
            slide,
            x + Inches(0.2),
            y + Inches(0.38),
            tc_w - Inches(0.4),
            Inches(0.25),
            text=spec,
            font_size=Pt(11),
            color=AI_NEUTRAL,
        )

    set_notes(slide, NOTES[2])


# ---------------------------------------------------------------------------
# Slide 3 — Radical Simplification
# ---------------------------------------------------------------------------
def slide_03_simplification(prs):
    slide = _blank_slide(prs)
    add_slide_header(
        slide,
        "Radical Simplification",
        subtitle="Spec 055 — Less is more: 1,016 files changed",
    )

    # Impact numbers row
    impact_y = Inches(2.2)
    impact_w = Inches(3.3)
    impact_h = Inches(1.3)
    impact_gap = Inches(0.3)
    impacts = [
        ("-21K", "Net Lines Removed", AI_SUCCESS),
        ("1,016", "Files Changed", AI_ACCENT),
        ("-74K", "Lines Deleted", AI_ERROR),
    ]
    total = len(impacts) * impact_w + (len(impacts) - 1) * impact_gap
    start_x = (SLIDE_W - total) // 2
    for i, (val, lbl, accent) in enumerate(impacts):
        x = start_x + i * (impact_w + impact_gap)
        add_card(
            slide,
            x,
            impact_y,
            impact_w,
            impact_h,
            fill_color=AI_CARD_DARK,
            border_color=AI_BORDER_DARK,
        )
        add_textbox(
            slide,
            x,
            impact_y + Inches(0.15),
            impact_w,
            Inches(0.6),
            text=val,
            font_name=FONT_TITLE,
            font_size=Pt(36),
            color=accent,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )
        add_textbox(
            slide,
            x,
            impact_y + Inches(0.8),
            impact_w,
            Inches(0.4),
            text=lbl,
            font_size=Pt(13),
            color=AI_NEUTRAL,
            alignment=PP_ALIGN.CENTER,
        )

    # What was removed / What was gained
    col_w = Inches(5.0)
    col_gap = Inches(0.8)
    total_cols = 2 * col_w + col_gap
    col_start_x = (SLIDE_W - total_cols) // 2
    col_y = Inches(3.9)

    # Removed column
    _add_feature_card(
        slide,
        col_start_x,
        col_y,
        col_w,
        Inches(2.6),
        "ELIMINATED",
        [
            "Programmatic pipeline generator (-2,175 LOC)",
            ".ai-engineering/skills/ & agents/ directories",
            "Stale spec save references",
            "Legacy naming conventions (ai: prefix)",
            "Redundant wrapper layers",
            "Suppression comments (noqa, nosec, etc.)",
        ],
        accent=AI_ERROR,
    )

    # Gained column
    _add_feature_card(
        slide,
        col_start_x + col_w + col_gap,
        col_y,
        col_w,
        Inches(2.6),
        "GAINED",
        [
            "IDE mirrors ARE the source of truth",
            "AI-driven pipeline generation via manifest",
            "Cleaner governance surface",
            "Simpler maintenance burden",
            "Consistent naming across all IDEs",
            "Validated by 1,696 passing tests",
        ],
        accent=AI_SUCCESS,
    )

    set_notes(slide, NOTES[3])


# ---------------------------------------------------------------------------
# Slide 4 — Multi-IDE Parity
# ---------------------------------------------------------------------------
def slide_04_multi_ide(prs):
    slide = _blank_slide(prs)
    add_slide_header(
        slide,
        "Multi-IDE Parity",
        subtitle="Specs 053, 062 + Copilot Hooks — Complete behavior in every IDE",
    )

    # Three IDE cards
    col_w = Inches(3.4)
    col_gap = Inches(0.25)
    total = 3 * col_w + 2 * col_gap
    start_x = (SLIDE_W - total) // 2
    col_y = Inches(2.0)

    # Claude Code
    _add_feature_card(
        slide,
        start_x,
        col_y,
        col_w,
        Inches(2.2),
        "Claude Code",
        [
            ".claude/skills/ai-*/SKILL.md",
            ".claude/agents/ai-*.md",
            "37 skills + 9 agents",
            "Hooks: pre/post tool, prompt submit",
            "Full telemetry + audit trail",
        ],
        accent=AI_ACCENT,
    )

    # GitHub Copilot
    _add_feature_card(
        slide,
        start_x + col_w + col_gap,
        col_y,
        col_w,
        Inches(2.2),
        "GitHub Copilot",
        [
            ".github/prompts/ai-*.prompt.md",
            ".github/agents/*.agent.md",
            "Handler concat with separators",
            "NEW: Copilot hooks parity",
            "Agent name detection in telemetry",
        ],
        accent=SEC_BLUE,
    )

    # Codex / Gemini
    _add_feature_card(
        slide,
        start_x + 2 * (col_w + col_gap),
        col_y,
        col_w,
        Inches(2.2),
        "Codex / Gemini CLI",
        [
            ".agents/skills/*/SKILL.md",
            ".agents/agents/ai-*.md",
            "AGENTS.md as entry point",
            "Same content, native format",
            "Handler files synced automatically",
        ],
        accent=SEC_PURPLE,
    )

    # Sync infrastructure
    sync_y = Inches(4.5)
    add_card(
        slide,
        LEFT_MARGIN,
        sync_y,
        CONTENT_W,
        Inches(1.4),
        fill_color=AI_CARD_DARK,
        border_color=AI_ACCENT,
    )
    add_textbox(
        slide,
        LEFT_MARGIN + Inches(0.2),
        sync_y + Inches(0.1),
        CONTENT_W - Inches(0.4),
        Inches(0.3),
        text="SYNC INFRASTRUCTURE (spec-062)",
        font_name=FONT_TITLE,
        font_size=Pt(14),
        color=AI_ACCENT,
        bold=True,
    )

    sync_items = [
        {
            "text": "  sync_command_mirrors.py: automated sync across .claude/, .github/, .agents/",
            "font_size": Pt(12),
            "color": AI_TEXT_PRIMARY,
        },
        {
            "text": "  Handler routing completeness: 90 regression tests",
            "font_size": Pt(12),
            "color": AI_TEXT_PRIMARY,
        },
        {
            "text": "  Template-prompt parity: 35 regression tests",
            "font_size": Pt(12),
            "color": AI_TEXT_PRIMARY,
        },
        {
            "text": "  Separators between handler sections in concatenated .prompt.md files",
            "font_size": Pt(12),
            "color": AI_TEXT_MUTED,
        },
    ]
    add_rich_textbox(
        slide,
        LEFT_MARGIN + Inches(0.2),
        sync_y + Inches(0.5),
        CONTENT_W - Inches(0.4),
        Inches(0.8),
        sync_items,
        line_spacing=Pt(8),
    )

    set_notes(slide, NOTES[4])


# ---------------------------------------------------------------------------
# Slide 5 — Pipeline v2 & AI-Driven Generation
# ---------------------------------------------------------------------------
def slide_05_pipeline(prs):
    slide = _blank_slide(prs)
    add_slide_header(
        slide,
        "Pipeline v2 & AI-Driven Generation",
        subtitle="Specs 054, 059 — From programmatic to intelligent",
    )

    # Timeline: 3 phases
    phase_y = Inches(2.2)
    phase_w = Inches(3.3)
    phase_h = Inches(3.0)
    phase_gap = Inches(0.3)
    total = 3 * phase_w + 2 * phase_gap
    start_x = (SLIDE_W - total) // 2

    # Phase 1: Pipeline v2
    _add_feature_card(
        slide,
        start_x,
        phase_y,
        phase_w,
        phase_h,
        "PHASE 1: CI Hardening (spec-054)",
        [
            "Generator rewrite for reliability",
            "GitHub Actions advanced patterns",
            "Azure Pipelines advanced patterns",
            "Telemetry hooks for skill tracking",
            "Guard events for governance",
            "Common templates across platforms",
            "Legacy naming cleanup (ai: prefix)",
        ],
        accent=SEC_BLUE,
    )

    # Arrow between phases
    arrow_x = start_x + phase_w + Inches(0.05)
    arrow = slide.shapes.add_shape(55, arrow_x, Inches(3.5), Inches(0.2), Inches(0.3))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = AI_ACCENT
    arrow.line.fill.background()

    # Phase 2: AI-Driven
    _add_feature_card(
        slide,
        start_x + phase_w + phase_gap,
        phase_y,
        phase_w,
        phase_h,
        "PHASE 2: AI-Driven (spec-059)",
        [
            "Eliminated programmatic generator",
            "Removed: cicd.py, injector.py, pipeline/",
            "Removed: 12 pipeline template files",
            "-2,175 lines of code deleted",
            "/ai-pipeline reads manifest.yml",
            "cicd.standards_url for team docs",
            "AI generates compliant pipelines",
        ],
        accent=AI_ACCENT,
    )

    # Arrow
    arrow_x2 = start_x + 2 * phase_w + phase_gap + Inches(0.05)
    arrow2 = slide.shapes.add_shape(55, arrow_x2, Inches(3.5), Inches(0.2), Inches(0.3))
    arrow2.fill.solid()
    arrow2.fill.fore_color.rgb = AI_ACCENT
    arrow2.line.fill.background()

    # Phase 3: Result
    _add_feature_card(
        slide,
        start_x + 2 * (phase_w + phase_gap),
        phase_y,
        phase_w,
        phase_h,
        "RESULT",
        [
            "Zero pipeline templates to maintain",
            "Team references their CI/CD docs",
            "AI generates platform-specific YAML",
            "Adapts to any CI provider",
            "No code changes needed for new stacks",
            "Simpler, more flexible, less brittle",
            "",
        ],
        accent=AI_SUCCESS,
    )

    # Bottom callout
    add_textbox(
        slide,
        LEFT_MARGIN,
        Inches(5.5),
        CONTENT_W,
        Inches(0.3),
        text=(
            "Philosophy: let AI do what AI does best — generate."
            " Maintain only the rules, not the templates."
        ),
        font_size=Pt(13),
        color=AI_NEUTRAL,
        alignment=PP_ALIGN.CENTER,
    )

    set_notes(slide, NOTES[5])


# ---------------------------------------------------------------------------
# Slide 6 — Enterprise & Delivery Features
# ---------------------------------------------------------------------------
def slide_06_enterprise(prs):
    slide = _blank_slide(prs)
    add_slide_header(
        slide,
        "Enterprise & Delivery Features",
        subtitle="Specs 056, 058 — Production readiness for teams",
    )

    # Two main feature areas side by side
    col_w = Inches(5.2)
    col_gap = Inches(0.4)
    total = 2 * col_w + col_gap
    start_x = (SLIDE_W - total) // 2
    col_y = Inches(2.0)

    # Enterprise Artifact Feed
    add_card(
        slide,
        start_x,
        col_y,
        col_w,
        Inches(2.5),
        fill_color=AI_CARD_DARK,
        border_color=AI_BORDER_DARK,
        left_accent_color=AI_WARNING,
    )
    add_textbox(
        slide,
        start_x + Inches(0.2),
        col_y + Inches(0.1),
        col_w - Inches(0.4),
        Inches(0.35),
        text="Enterprise Artifact Feed (spec-056)",
        font_name=FONT_TITLE,
        font_size=Pt(16),
        color=AI_WARNING,
        bold=True,
    )
    feed_lines = [
        {"text": "  Manifest.yml reorganization", "font_size": Pt(12), "color": AI_TEXT_PRIMARY},
        {
            "text": "  Enterprise artifact feed configuration",
            "font_size": Pt(12),
            "color": AI_TEXT_PRIMARY,
        },
        {
            "text": "  Clean separation of team vs framework config",
            "font_size": Pt(12),
            "color": AI_TEXT_PRIMARY,
        },
        {
            "text": "  Consistent YAML structure across environments",
            "font_size": Pt(12),
            "color": AI_TEXT_MUTED,
        },
    ]
    add_rich_textbox(
        slide,
        start_x + Inches(0.2),
        col_y + Inches(0.55),
        col_w - Inches(0.4),
        Inches(1.8),
        feed_lines,
        line_spacing=Pt(10),
    )

    # Watch & Fix Loop
    wf_x = start_x + col_w + col_gap
    add_card(
        slide,
        wf_x,
        col_y,
        col_w,
        Inches(2.5),
        fill_color=AI_CARD_DARK,
        border_color=AI_BORDER_DARK,
        left_accent_color=SEC_BLUE,
    )
    add_textbox(
        slide,
        wf_x + Inches(0.2),
        col_y + Inches(0.1),
        col_w - Inches(0.4),
        Inches(0.35),
        text="Watch & Fix Loop (spec-058)",
        font_name=FONT_TITLE,
        font_size=Pt(16),
        color=SEC_BLUE,
        bold=True,
    )
    wf_lines = [
        {"text": "  Auto-retry for /ai-pr step 14", "font_size": Pt(12), "color": AI_TEXT_PRIMARY},
        {
            "text": "  Handles transient CI failures gracefully",
            "font_size": Pt(12),
            "color": AI_TEXT_PRIMARY,
        },
        {
            "text": "  Monitors PR checks and re-triggers fixes",
            "font_size": Pt(12),
            "color": AI_TEXT_PRIMARY,
        },
        {
            "text": "  1,232 lines of delivery infrastructure",
            "font_size": Pt(12),
            "color": AI_TEXT_MUTED,
        },
    ]
    add_rich_textbox(
        slide,
        wf_x + Inches(0.2),
        col_y + Inches(0.55),
        col_w - Inches(0.4),
        Inches(1.8),
        wf_lines,
        line_spacing=Pt(10),
    )

    # Skill frontmatter evolution
    fm_y = Inches(4.8)
    add_card(
        slide,
        LEFT_MARGIN,
        fm_y,
        CONTENT_W,
        Inches(1.3),
        fill_color=AI_CARD_DARK,
        border_color=AI_BORDER_DARK,
        left_accent_color=SEC_PURPLE,
    )
    add_textbox(
        slide,
        LEFT_MARGIN + Inches(0.2),
        fm_y + Inches(0.1),
        CONTENT_W - Inches(0.4),
        Inches(0.3),
        text="Skill Framework Evolution",
        font_name=FONT_TITLE,
        font_size=Pt(14),
        color=SEC_PURPLE,
        bold=True,
    )
    fm_lines = [
        {
            "text": "  Added effort levels to all 37 skills: max (11), high (15), medium (11)",
            "font_size": Pt(12),
            "color": AI_TEXT_PRIMARY,
        },
        {
            "text": (
                "  Effort determines cognitive weight"
                " — used for compute allocation and agent selection"
            ),
            "font_size": Pt(12),
            "color": AI_TEXT_PRIMARY,
        },
        {
            "text": "  221 files updated across all IDE mirrors to include effort frontmatter",
            "font_size": Pt(12),
            "color": AI_TEXT_MUTED,
        },
    ]
    add_rich_textbox(
        slide,
        LEFT_MARGIN + Inches(0.2),
        fm_y + Inches(0.45),
        CONTENT_W - Inches(0.4),
        Inches(0.7),
        fm_lines,
        line_spacing=Pt(8),
    )

    set_notes(slide, NOTES[6])


# ---------------------------------------------------------------------------
# Slide 7 — CI Reliability
# ---------------------------------------------------------------------------
def slide_07_ci_reliability(prs):
    slide = _blank_slide(prs)
    add_slide_header(
        slide,
        "CI Reliability & Test Infrastructure",
        subtitle="Spec 060 + community fixes — Trust in the pipeline",
    )

    # Fixes table
    tbl, _shape = add_styled_table(
        slide,
        LEFT_MARGIN,
        Inches(2.2),
        CONTENT_W,
        Inches(3.5),
        8,
        3,
        header_fill=AI_ACCENT,
    )

    headers = ["Area", "Change", "Impact"]
    data = [
        ("CI False Positives", "Eliminated spurious test failures (spec-060)", "550+ LOC of fixes"),
        ("Test Mapping", "Fixed test_scope fallback + CI mapping", "Correct test selection"),
        ("Dependabot", "Grouped PRs, exempted from trailer checks", "Less CI noise"),
        ("Governance", "Restored validation to PASS state", "-1,829 LOC cleanup"),
        ("Test Rewrites", "test_doctor: 63 patches to 1 fixture", "Maintainable tests"),
        ("Suppression", "Removed all noqa/nosec/type:ignore", "Clean codebase"),
        ("Documentation", "Updated skill/agent counts, removed stale refs", "Accurate docs"),
    ]

    for ci, h in enumerate(headers):
        _style_cell(
            tbl.cell(0, ci),
            h,
            bold=True,
            color=AI_BG_DARK,
            font_size=Pt(13),
            alignment=PP_ALIGN.CENTER,
        )

    for ri, (area, change, impact) in enumerate(data, start=1):
        bg = AI_CARD_DARK if ri % 2 == 0 else AI_BG_DARK
        _style_cell(
            tbl.cell(ri, 0), area, font_size=Pt(11), fill_color=bg, color=AI_ACCENT, bold=True
        )
        _style_cell(tbl.cell(ri, 1), change, font_size=Pt(11), fill_color=bg, color=AI_TEXT_PRIMARY)
        _style_cell(tbl.cell(ri, 2), impact, font_size=Pt(11), fill_color=bg, color=AI_TEXT_MUTED)

    # Bottom note
    add_textbox(
        slide,
        LEFT_MARGIN,
        Inches(6.0),
        CONTENT_W,
        Inches(0.3),
        text="Every fix validated by full test suite — no regression introduced",
        font_size=Pt(12),
        color=AI_NEUTRAL,
        alignment=PP_ALIGN.CENTER,
    )

    set_notes(slide, NOTES[7])


# ---------------------------------------------------------------------------
# Slide 8 — Quality Metrics
# ---------------------------------------------------------------------------
def slide_08_quality(prs):
    slide = _blank_slide(prs)
    add_slide_header(slide, "Quality Metrics", subtitle="Enforced — not aspirational")

    # KPI row
    kpi_y = Inches(2.2)
    kpi_w = Inches(2.0)
    kpi_h = Inches(1.3)
    kpi_gap = Inches(0.25)
    kpis = [
        ("~2K", "Tests", AI_ACCENT),
        ("37", "Skills", SEC_BLUE),
        ("9", "Agents", SEC_PURPLE),
        ("4", "IDEs", AI_WARNING),
        ("100%", "Gate Pass", AI_SUCCESS),
    ]
    total = len(kpis) * kpi_w + (len(kpis) - 1) * kpi_gap
    start_x = (SLIDE_W - total) // 2

    for i, (val, lbl, accent) in enumerate(kpis):
        x = start_x + i * (kpi_w + kpi_gap)
        _add_kpi_card(slide, x, kpi_y, kpi_w, kpi_h, val, lbl, accent=accent)

    # Sprint contribution summary
    tbl, _shape = add_styled_table(
        slide,
        LEFT_MARGIN,
        Inches(4.0),
        CONTENT_W,
        Inches(2.2),
        5,
        4,
        header_fill=AI_ACCENT,
    )

    headers = ["Metric", "Last Sprint", "This Sprint", "Delta"]
    data = [
        ("Skills", "40", "37", "-3 (consolidated)"),
        ("Agents", "10", "9", "-1 (consolidated)"),
        ("Pipeline Code", "2,175 LOC", "0 LOC", "-100% (AI-driven)"),
        ("Tests", "1,696", "~1,968", "+16%"),
    ]

    for ci, h in enumerate(headers):
        _style_cell(
            tbl.cell(0, ci),
            h,
            bold=True,
            color=AI_BG_DARK,
            font_size=Pt(13),
            alignment=PP_ALIGN.CENTER,
        )

    for ri, (metric, last, this, delta) in enumerate(data, start=1):
        bg = AI_CARD_DARK if ri % 2 == 0 else AI_BG_DARK
        _style_cell(tbl.cell(ri, 0), metric, font_size=Pt(12), fill_color=bg, color=AI_TEXT_PRIMARY)
        _style_cell(
            tbl.cell(ri, 1),
            last,
            font_size=Pt(12),
            fill_color=bg,
            color=AI_TEXT_MUTED,
            alignment=PP_ALIGN.CENTER,
        )
        _style_cell(
            tbl.cell(ri, 2),
            this,
            font_size=Pt(12),
            fill_color=bg,
            color=AI_WHITE,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )
        delta_color = (
            AI_SUCCESS if "+" in delta or "-100" in delta or "consolidated" in delta else AI_WARNING
        )
        _style_cell(
            tbl.cell(ri, 3),
            delta,
            font_size=Pt(12),
            fill_color=bg,
            color=delta_color,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )

    set_notes(slide, NOTES[8])


# ---------------------------------------------------------------------------
# Slide 9 — Risks & Next Sprint
# ---------------------------------------------------------------------------
def slide_09_risks_next(prs):
    slide = _blank_slide(prs)
    add_slide_header(slide, "Risks & Next Sprint")

    # Risks section
    risk_x = LEFT_MARGIN
    risk_y = Inches(2.0)
    risk_w = Inches(5.2)

    add_textbox(
        slide,
        risk_x,
        risk_y,
        risk_w,
        Inches(0.35),
        text="ACTIVE RISKS",
        font_name=FONT_TITLE,
        font_size=Pt(16),
        color=AI_ERROR,
        bold=True,
    )

    risks = [
        (
            "MEDIUM",
            "Uncommitted new skills in progress",
            "autopilot, eval, media, slides, video-editing — in development on main",
        ),
        (
            "LOW",
            "Cross-IDE validation pending",
            "Full end-to-end verification across all 4 IDEs not yet automated",
        ),
        (
            "LOW",
            "Template-project drift monitoring",
            "New handler files need ongoing sync validation",
        ),
    ]
    for i, (severity, title, desc) in enumerate(risks):
        ry = risk_y + Inches(0.5) + i * Inches(0.85)
        sev_color = AI_WARNING if severity == "MEDIUM" else AI_ACCENT
        add_card(
            slide,
            risk_x,
            ry,
            risk_w,
            Inches(0.75),
            fill_color=AI_CARD_DARK,
            border_color=AI_BORDER_DARK,
            left_accent_color=sev_color,
        )
        add_textbox(
            slide,
            risk_x + Inches(0.2),
            ry + Inches(0.08),
            Inches(1.0),
            Inches(0.25),
            text=severity,
            font_name=FONT_TITLE,
            font_size=Pt(10),
            color=sev_color,
            bold=True,
        )
        add_textbox(
            slide,
            risk_x + Inches(1.2),
            ry + Inches(0.08),
            risk_w - Inches(1.4),
            Inches(0.25),
            text=title,
            font_size=Pt(12),
            color=AI_TEXT_PRIMARY,
            bold=True,
        )
        add_textbox(
            slide,
            risk_x + Inches(0.2),
            ry + Inches(0.38),
            risk_w - Inches(0.4),
            Inches(0.3),
            text=desc,
            font_size=Pt(10),
            color=AI_TEXT_MUTED,
        )

    # Next Sprint section
    next_x = risk_x + risk_w + Inches(0.4)
    next_y = Inches(2.0)
    next_w = Inches(5.2)

    add_textbox(
        slide,
        next_x,
        next_y,
        next_w,
        Inches(0.35),
        text="NEXT SPRINT PRIORITIES",
        font_name=FONT_TITLE,
        font_size=Pt(16),
        color=AI_SUCCESS,
        bold=True,
    )

    nexts = [
        ("New Skills Completion", "Land autopilot, eval, media, slides, video-editing", AI_ACCENT),
        (
            "Cross-IDE E2E Validation",
            "Automated verification across Claude, Copilot, Gemini, Codex",
            SEC_BLUE,
        ),
        (
            "Documentation Site",
            "Public docs site for framework onboarding and reference",
            SEC_PURPLE,
        ),
        (
            "Instinct System",
            "Session learning and pattern evolution across conversations",
            AI_WARNING,
        ),
    ]
    for i, (title, desc, accent) in enumerate(nexts):
        ny = next_y + Inches(0.5) + i * Inches(0.85)
        add_card(
            slide,
            next_x,
            ny,
            next_w,
            Inches(0.75),
            fill_color=AI_CARD_DARK,
            border_color=AI_BORDER_DARK,
            left_accent_color=accent,
        )
        add_textbox(
            slide,
            next_x + Inches(0.2),
            ny + Inches(0.08),
            next_w - Inches(0.4),
            Inches(0.25),
            text=title,
            font_size=Pt(13),
            color=accent,
            bold=True,
        )
        add_textbox(
            slide,
            next_x + Inches(0.2),
            ny + Inches(0.38),
            next_w - Inches(0.4),
            Inches(0.3),
            text=desc,
            font_size=Pt(11),
            color=AI_TEXT_MUTED,
        )

    set_notes(slide, NOTES[9])


# ---------------------------------------------------------------------------
# Slide 10 — Q&A
# ---------------------------------------------------------------------------
def slide_10_qa(prs):
    slide = _blank_slide(prs)

    add_accent_bar(slide, Inches(0), Inches(0.4), SLIDE_W, height=Pt(4))

    add_textbox(
        slide,
        LEFT_MARGIN,
        Inches(2.5),
        CONTENT_W,
        Inches(1.0),
        text="Thank You",
        font_name=FONT_TITLE,
        font_size=Pt(52),
        color=AI_WHITE,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )

    bar_w = Inches(3)
    bar_left = (SLIDE_W - bar_w) // 2
    add_accent_bar(slide, bar_left, Inches(3.5), bar_w, height=Pt(3))

    add_textbox(
        slide,
        LEFT_MARGIN,
        Inches(3.8),
        CONTENT_W,
        Inches(0.5),
        text="Questions?",
        font_size=Pt(24),
        color=AI_ACCENT,
        alignment=PP_ALIGN.CENTER,
    )

    # Summary badges
    badge_y = Inches(5.0)
    badge_w = Inches(2.2)
    badge_h = Inches(0.5)
    gap = Inches(0.3)
    badges = [
        "Simplification",
        "Multi-IDE",
        "AI Pipelines",
        "Enterprise",
    ]
    total_w = len(badges) * badge_w + (len(badges) - 1) * gap
    start_x = (SLIDE_W - total_w) // 2
    for i, label in enumerate(badges):
        x = start_x + i * (badge_w + gap)
        add_rounded_rect(
            slide, x, badge_y, badge_w, badge_h, fill_color=AI_CARD_DARK, border_color=AI_ACCENT
        )
        add_textbox(
            slide,
            x,
            badge_y,
            badge_w,
            badge_h,
            text=label,
            font_size=Pt(12),
            color=AI_ACCENT,
            bold=True,
            alignment=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )

    add_textbox(
        slide,
        LEFT_MARGIN,
        Inches(5.8),
        CONTENT_W,
        Inches(0.4),
        text="Sprint Mar 16 - Mar 24, 2026  |  ai-engineering",
        font_size=Pt(14),
        color=AI_NEUTRAL,
        alignment=PP_ALIGN.CENTER,
    )

    add_accent_bar(slide, Inches(0), Inches(7.0), SLIDE_W, height=Pt(4))
    set_notes(slide, NOTES[10])


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_01_title(prs)
    slide_02_overview(prs)
    slide_03_simplification(prs)
    slide_04_multi_ide(prs)
    slide_05_pipeline(prs)
    slide_06_enterprise(prs)
    slide_07_ci_reliability(prs)
    slide_08_quality(prs)
    slide_09_risks_next(prs)
    slide_10_qa(prs)

    out = Path(__file__).parent / "sprint-review-2026-03-24.pptx"
    prs.save(str(out))
    print(f"Generated: {out}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
