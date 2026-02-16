"""Generate ai-engineering board presentation (.pptx) with ai-engineering dark-mode brand identity.

Output: docs/presentation/ai-engineering-board.pptx
Brand: ai-engineering dark mode — dark backgrounds, teal accents, technical tone.
All shapes are native PowerPoint objects (fully editable).
"""

from __future__ import annotations

import contextlib
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# ai-engineering Dark-Mode Colour Palette  (from docs/design.pen variables)
# ---------------------------------------------------------------------------
AI_BG_DARK = RGBColor(0x0B, 0x11, 0x20)  # $bg-dark / $primary-dark
AI_WHITE = RGBColor(0xFF, 0xFF, 0xFF)  # $white — headings only
AI_TEXT_PRIMARY = RGBColor(0xE2, 0xE8, 0xF0)  # $text-primary — body text

# Accent / structural
AI_ACCENT = RGBColor(0x00, 0xD4, 0xAA)  # $accent
AI_PRIMARY = RGBColor(0x1E, 0x3A, 0x5F)  # $primary
AI_PRIMARY_LIGHT = RGBColor(0x2A, 0x4F, 0x7A)  # $primary-light
AI_ERROR = RGBColor(0xEF, 0x44, 0x44)
AI_SUCCESS = RGBColor(0x10, 0xB9, 0x81)
AI_WARNING = RGBColor(0xF5, 0x9E, 0x0B)

# Text / neutral scale
AI_TEXT_LIGHT = RGBColor(0xF8, 0xFA, 0xFB)  # $light
AI_NEUTRAL = RGBColor(0x64, 0x74, 0x8B)  # $neutral
AI_TEXT_MUTED = RGBColor(0x94, 0xA3, 0xB8)  # $text-muted

# Borders and cards
AI_BORDER_DARK = RGBColor(0x1A, 0x2A, 0x40)  # $border-dark
AI_BORDER_LIGHT = RGBColor(0xE2, 0xE8, 0xF0)  # $border-light
AI_CARD_DARK = RGBColor(0x1E, 0x29, 0x3B)  # $card-bg

# Secondary palette — for charts/diagrams differentiation
SEC_BLUE_DARK = RGBColor(0x1A, 0x3A, 0x5C)
SEC_BLUE = RGBColor(0x2E, 0x6B, 0xA4)
SEC_BLUE_LIGHT = RGBColor(0x7A, 0xB5, 0xD6)
SEC_BLUE_PALE = RGBColor(0x0B, 0x1E, 0x3A)

SEC_TEAL_DARK = RGBColor(0x00, 0x7A, 0x62)
SEC_TEAL = RGBColor(0x00, 0xD4, 0xAA)
SEC_TEAL_LIGHT = RGBColor(0x5C, 0xE8, 0xCC)
SEC_TEAL_PALE = RGBColor(0x0B, 0x2A, 0x22)

SEC_PURPLE_DARK = RGBColor(0x4A, 0x1A, 0x6B)
SEC_PURPLE = RGBColor(0x7B, 0x3F, 0xA0)
SEC_PURPLE_LIGHT = RGBColor(0xB3, 0x8B, 0xCF)
SEC_PURPLE_PALE = RGBColor(0x1A, 0x0B, 0x2E)

# ---------------------------------------------------------------------------
# Typography
# ---------------------------------------------------------------------------
FONT_TITLE = "JetBrains Mono"
FONT_TITLE_FALLBACK = "monospace"
FONT_BODY = "Inter"
FONT_BODY_FALLBACK = "sans-serif"

# ---------------------------------------------------------------------------
# Slide dimensions (16:9 widescreen)
# ---------------------------------------------------------------------------
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

LEFT_MARGIN = Inches(1.2)
RIGHT_MARGIN = Inches(1.2)
CONTENT_W = SLIDE_W - LEFT_MARGIN - RIGHT_MARGIN

# Default sizes (avoid B008 — function calls in arg defaults)
_PT_3 = Pt(3)
_PT_14 = Pt(14)
_PT_20 = Pt(20)
_PT_1 = Pt(1)
_PT_12 = Pt(12)
_IN_008 = Inches(0.08)

# ---------------------------------------------------------------------------
# Speaker notes — English talking points from speech-script.md
# ---------------------------------------------------------------------------
NOTES: dict[int, str] = {
    1: (
        "Good morning. Today I present ai-engineering — a governance framework "
        "for AI-assisted development.\n\n"
        "It is not a platform to buy. It is not a pipeline to manage. It is a "
        "content framework — Markdown, YAML, JSON, Bash — that turns AI "
        "assistance into governed delivery.\n\n"
        "Open source, MIT license, Python 3.11+, compatible with any operating system."
    ),
    2: (
        "To understand why we need governance, let's see how we got here.\n\n"
        "In 2022 we had code completion — Copilot, TabNine — suggesting line by line. "
        "In 2023 came chat-in-IDE: conversations with AI inside the editor. "
        "In 2024, agentic coding — tools like Claude Code or Devin that execute complete "
        "tasks autonomously.\n\n"
        "In 2025 the multi-agent ecosystem exploded: MCP — Model Context Protocol — for "
        "tool integration, and A2A — Agent-to-Agent — for agent coordination.\n\n"
        "But here is the critical point: all this capability has developed without a "
        "governance layer. More power without more control.\n\n"
        "Skills: reusable procedures in Markdown.\n"
        "Agents: specialized personas — behavior contracts.\n"
        "MCP: Model Context Protocol — standard for external tool connection.\n"
        "A2A: Agent-to-Agent — coordination between multiple agents.\n\n"
        "ai-engineering does NOT use MCP or A2A directly — but it positions itself as the "
        "missing governance layer."
    ),
    3: (
        "This is what happens when you use AI to code without governance.\n\n"
        "Four AI agents, same codebase, zero coordination. The result:\n"
        "- Secrets in commits — AI generates code with credentials.\n"
        "- Quality gates bypassed — AI skips tests, ignores linting.\n"
        "- Architectural drift — each agent makes different decisions.\n"
        "- Repeated decisions — AI asks the same thing every session.\n\n"
        "For the boards:\n"
        "- Compliance gaps: no audit trail.\n"
        "- Security exposure: secrets in git.\n"
        "- Quality degradation: no thresholds.\n"
        "- Knowledge loss: decisions that are lost."
    ),
    4: (
        "We did not arrive at ai-engineering on the first try. We explored four frameworks:\n\n"
        "- SpecKit: good spec management, but no enforcement.\n"
        "- BMAD Method: excellent orchestration, but heavyweight.\n"
        "- GSD: pragmatic, but no governance.\n"
        "- OpenSpec: standard aspirations, but no practical tooling.\n\n"
        "After 5 iterations: content-first governance, simple to adopt, "
        "strict for enforcement, flexible to scale."
    ),
    5: (
        "What is ai-engineering exactly?\n\n"
        "It is NOT a platform. It is NOT a CI/CD pipeline. It is a content framework.\n\n"
        "Five subdirectories:\n"
        "- standards/ — framework and team rules.\n"
        "- skills/ — 31 reusable procedures in 7 categories.\n"
        "- agents/ — 8 specialized personas.\n"
        "- context/ — delivery specs, product contracts.\n"
        "- state/ — decision store, audit log, manifests.\n\n"
        "Minimal CLI: ai-eng install, update, doctor, validate.\n\n"
        "Works with Claude Code (37 slash commands), GitHub Copilot, OpenAI Codex.\n"
        "A constitution for the AI-assisted repository."
    ),
    6: (
        "The developer experience in under 5 minutes:\n\n"
        "ai-eng install . — creates the governance root, "
        "configures git hooks, generates state files.\n\n"
        "Every commit goes through quality gates:\n"
        "- Pre-commit: ruff format/lint, gitleaks.\n"
        "- Commit-msg: valid format, branch protection.\n"
        "- Pre-push: semgrep SAST/OWASP, pip-audit CVEs, pytest, ty type checking.\n\n"
        "If any gate fails: the push is blocked. No bypass.\n\n"
        "Your next commit after install is already governed."
    ),
    7: (
        "Four clear boundaries:\n\n"
        "1. Framework-managed: standards, skills, agents. Updatable with ai-eng update.\n"
        "2. Team-managed: standards/team/. NEVER overwritten by updates.\n"
        "3. Project-managed: specs, contracts. NEVER overwritten.\n"
        "4. System-managed: runtime state files.\n\n"
        "Ownership boundaries are non-negotiable."
    ),
    8: (
        "31 skills organized in 7 categories:\n"
        "- Workflows (4): commit, PR, acho, pre-implementation.\n"
        "- Dev (6): debug, refactor, code review, test strategy, migration, deps.\n"
        "- Review (3): architecture, performance, security.\n"
        "- Docs (4): changelog, explain, writer, prompt design.\n"
        "- Govern (9): create/delete specs, skills, agents + risk lifecycle.\n"
        "- Quality (3): audit code, audit report, install check.\n"
        "- Utils (3): git helpers, platform detection, Python patterns.\n\n"
        "They are NOT code. They are behavior specifications. Any AI agent "
        "reads and executes them."
    ),
    9: (
        "8 specialized agents:\n"
        "- Principal Engineer: senior code review.\n"
        "- Architect: architectural analysis.\n"
        "- Security Reviewer: security assessment.\n"
        "- Debugger: systematic diagnosis.\n"
        "- Quality Auditor: quality gate enforcement.\n"
        "- Verify App: end-to-end verification.\n"
        "- Codebase Mapper: structure mapping.\n"
        "- Code Simplifier: complexity reduction.\n\n"
        "Each agent has: Identity, Capabilities, Activation rules, Behavior protocol, "
        "Referenced Skills, Output Contract, and Boundaries."
    ),
    10: (
        "Every non-trivial change follows a 4-document cycle:\n"
        "- spec.md — the WHAT.\n"
        "- plan.md — the HOW.\n"
        "- tasks.md — the DO.\n"
        "- done.md — the DONE.\n\n"
        "It is NOT bureaucracy. It is AI session recovery. Any agent can resume "
        "any spec at any point.\n\n"
        "Multi-agent execution is parallel. Each phase passes through a phase gate."
    ),
    11: (
        "Institutional memory lives in 5 state files:\n"
        "- install-manifest.json: what was installed, when.\n"
        "- ownership-map.json: who owns each path.\n"
        "- sources.lock.json: remote skills with verifiable checksums.\n"
        "- decision-store.json: 10 real decisions with SHA-256 context hash.\n"
        "- audit-log.ndjson: 183 recorded events.\n\n"
        "Decision continuity: Agent A decides in session 1, "
        "Agent B does not ask again in session 5."
    ),
    12: (
        "3 mandatory stages:\n"
        "- Pre-commit: ruff format, ruff lint, gitleaks.\n"
        "- Commit-msg: valid format, branch protection.\n"
        "- Pre-push: semgrep, pip-audit, pytest, ty.\n\n"
        "Thresholds: Coverage >=80%, Duplication <=3%, CC <=10, CogC <=15.\n\n"
        "Structured risk acceptance: Critical 15d, High 30d, Medium 60d, Low 90d. "
        "Maximum 2 renewals."
    ),
    13: (
        "The value depends on who is looking:\n"
        "- Engineers: 37 slash commands, quality gates before push.\n"
        "- Governance/Compliance: audit-log with 183 traceable events, "
        "decision store with 10 decisions.\n"
        "- Security/AppSec: gitleaks + semgrep + pip-audit on every push.\n"
        "- Quality/DevEx: Sonar-like quality gates WITHOUT a SonarQube server.\n"
        "- Architecture: Standards with layering, ownership boundaries."
    ),
    14: (
        "Direct business case:\n"
        "- Risk reduction: 0 ungated operations. 100% gate execution.\n"
        "- Cost avoidance: 0 SonarQube licenses. MIT open source.\n"
        "- Consistency: same framework across all repos.\n"
        "- Time savings: <5 min from install to first governed commit.\n"
        "- Compliance readiness: audit log, risk acceptance, decision store with SHA-256."
    ),
    15: (
        "No vendor lock-in.\n\n"
        "Claude Code: CLAUDE.md + 37 slash commands.\n"
        "Copilot: copilot-instructions.md + prompts in .github/copilot/.\n"
        "Codex: codex.md.\n"
        "Terminal: CLI directly.\n\n"
        "You change providers, you keep the governance."
    ),
    16: (
        "Why not just an instructions file?\n\n"
        "Simple instructions work for individual tasks. They break at scale:\n"
        "- No enforcement, no state, no ownership, no audit trail.\n"
        "- No security scanning, no risk management, no delivery lifecycle.\n\n"
        "Simple instructions: 0/8 capabilities enforced. ai-engineering: 8/8.\n\n"
        "Instructions tell AI what to do. ai-engineering ensures it actually does it."
    ),
    17: (
        "Against the alternatives:\n"
        "- vs SpecKit: manages specs — ai-engineering manages the full governed cycle.\n"
        "- vs BMAD: strong multi-agent, but heavyweight.\n"
        "- vs GSD: pragmatic, no governance backbone.\n"
        "- vs OpenSpec: aspirational standard, not installable today.\n\n"
        "Five differentiators:\n"
        "1. Content-first\n"
        "2. Non-bypassable enforcement\n"
        "3. Risk lifecycle\n"
        "4. Cross-IDE day one\n"
        "5. Ownership model"
    ),
    18: (
        "The ask: approve ai-engineering as the governance standard for AI-assisted "
        "development in the organization.\n\n"
        "Pilot plan: 2-3 repositories next quarter.\n"
        "Metrics: gate execution rate, time to governed commit, security catch rate, "
        "decision reuse rate.\n\n"
        "Roadmap:\n"
        "- Phase 1: GitHub + Python + Claude/Copilot/Codex.\n"
        "- Phase 2: Azure DevOps + more stacks + signature verification.\n"
        "- Phase 3: Multi-agent orchestration + docs site.\n\n"
        "Investment: $0 license cost. MIT open source.\n"
        "Next step: Approve the pilot scope. The framework is ready to install today.\n\n"
        "Thank you. Questions?"
    ),
}


# ---------------------------------------------------------------------------
# Helper utilities
# ---------------------------------------------------------------------------
def _font(
    run,
    *,
    name: str = FONT_BODY,
    size: Pt = _PT_14,
    color: RGBColor = AI_TEXT_LIGHT,
    bold: bool = False,
):
    """Apply font styling to a text run."""
    run.font.name = name
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    font_name: str = FONT_BODY,
    font_size: Pt = _PT_14,
    color: RGBColor = AI_TEXT_LIGHT,
    bold: bool = False,
    alignment: PP_ALIGN = PP_ALIGN.LEFT,
    word_wrap: bool = True,
    anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
):
    """Add a text box with single-style text and return the shape."""
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = word_wrap
    tf.auto_size = None
    with contextlib.suppress(Exception):
        tf.paragraphs[0].alignment = alignment
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    _font(run, name=font_name, size=font_size, color=color, bold=bold)
    # Vertical anchor
    txbox.text_frame.paragraphs[0].space_before = Pt(0)
    txbox.text_frame.paragraphs[0].space_after = Pt(0)
    with contextlib.suppress(Exception):
        txbox.text_frame._txBody.bodyPr.set(
            "anchor",
            {
                MSO_ANCHOR.TOP: "t",
                MSO_ANCHOR.MIDDLE: "ctr",
                MSO_ANCHOR.BOTTOM: "b",
            }.get(anchor, "t"),
        )
    return txbox


def add_rich_textbox(
    slide, left, top, width, height, lines, *, alignment=PP_ALIGN.LEFT, line_spacing=_PT_20
):
    """Add a text box with multiple styled lines.

    *lines* is a list of dicts with keys: text, font_name, font_size, color, bold.
    Each dict becomes a separate paragraph.
    """
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = alignment
        if line_spacing:
            p.space_after = line_spacing
        run = p.add_run()
        run.text = ln.get("text", "")
        _font(
            run,
            name=ln.get("font_name", FONT_BODY),
            size=ln.get("font_size", Pt(14)),
            color=ln.get("color", AI_TEXT_LIGHT),
            bold=ln.get("bold", False),
        )
    return txbox


def add_accent_bar(slide, left, top, width, *, height=_PT_3, color=AI_ACCENT):
    """Thin horizontal accent bar — teal signature brand element."""
    bar = slide.shapes.add_shape(1, left, top, width, height)  # 1 = MSO_SHAPE.RECTANGLE
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar


def add_card(
    slide,
    left,
    top,
    width,
    height,
    *,
    fill_color=AI_CARD_DARK,
    border_color=None,
    border_width=_PT_1,
    left_accent_color=None,
    left_accent_width=_IN_008,
):
    """Add a card rectangle (dark bg, optional border / left accent)."""
    card = slide.shapes.add_shape(1, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = fill_color
    if border_color:
        card.line.color.rgb = border_color
        card.line.width = border_width
    else:
        card.line.fill.background()
    # Left-accent bar overlaid on the card
    if left_accent_color:
        accent = slide.shapes.add_shape(1, left, top, left_accent_width, height)
        accent.fill.solid()
        accent.fill.fore_color.rgb = left_accent_color
        accent.line.fill.background()
    return card


def add_circle(slide, left, top, diameter, *, fill_color=AI_ACCENT, border_color=None):
    """Add a circle (oval shape with equal w/h)."""
    circ = slide.shapes.add_shape(9, left, top, diameter, diameter)  # 9 = MSO_SHAPE.OVAL
    circ.fill.solid()
    circ.fill.fore_color.rgb = fill_color
    if border_color:
        circ.line.color.rgb = border_color
        circ.line.width = Pt(1)
    else:
        circ.line.fill.background()
    return circ


def add_arrow_right(slide, left, top, width, height, *, color=AI_ACCENT):
    """Add a right-pointing arrow shape."""
    # 55 = MSO_SHAPE.RIGHT_ARROW
    arrow = slide.shapes.add_shape(55, left, top, width, height)
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = color
    arrow.line.fill.background()
    return arrow


def add_chevron(slide, left, top, width, height, *, color=AI_ACCENT):
    """Add a chevron (notched right arrow)."""
    # 94 = MSO_SHAPE.CHEVRON
    chev = slide.shapes.add_shape(94, left, top, width, height)
    chev.fill.solid()
    chev.fill.fore_color.rgb = color
    chev.line.fill.background()
    return chev


def add_rect(
    slide,
    left,
    top,
    width,
    height,
    *,
    fill_color=AI_CARD_DARK,
    border_color=AI_TEXT_MUTED,
    border_width=_PT_1,
):
    """Add a rectangle with optional fill and border."""
    rect = slide.shapes.add_shape(1, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill_color
    if border_color:
        rect.line.color.rgb = border_color
        rect.line.width = border_width
    else:
        rect.line.fill.background()
    return rect


def add_rounded_rect(
    slide, left, top, width, height, *, fill_color=AI_CARD_DARK, border_color=AI_TEXT_MUTED
):
    """Rounded rectangle."""
    rr = slide.shapes.add_shape(5, left, top, width, height)  # 5 = MSO_SHAPE.ROUNDED_RECTANGLE
    rr.fill.solid()
    rr.fill.fore_color.rgb = fill_color
    if border_color:
        rr.line.color.rgb = border_color
        rr.line.width = Pt(1)
    else:
        rr.line.fill.background()
    return rr


def set_notes(slide, text: str):
    """Set speaker notes for a slide."""
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


def add_slide_header(slide, title: str, *, subtitle: str | None = None):
    """Standard slide header: title + accent bar + optional subtitle."""
    add_accent_bar(slide, LEFT_MARGIN, Inches(0.6), CONTENT_W)
    add_textbox(
        slide,
        LEFT_MARGIN,
        Inches(0.75),
        CONTENT_W,
        Inches(0.7),
        text=title,
        font_name=FONT_TITLE,
        font_size=Pt(34),
        color=AI_WHITE,
        bold=True,
    )
    if subtitle:
        add_textbox(
            slide,
            LEFT_MARGIN,
            Inches(1.35),
            CONTENT_W,
            Inches(0.4),
            text=subtitle,
            font_name=FONT_BODY,
            font_size=Pt(18),
            color=AI_NEUTRAL,
        )


def _blank_slide(prs):
    """Add a blank slide with dark background and return it."""
    layout = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(layout)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = AI_BG_DARK
    return slide


# ---------------------------------------------------------------------------
# Styled table helper
# ---------------------------------------------------------------------------
def add_styled_table(
    slide,
    left,
    top,
    width,
    height,
    rows,
    cols,
    *,
    header_fill=AI_ACCENT,
    header_font_color=AI_TEXT_PRIMARY,
    alt_row_fill=AI_CARD_DARK,
):
    """Add a table shape and return (table, shape) for further customisation."""
    shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    tbl = shape.table
    # Style header row
    for ci in range(cols):
        cell = tbl.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
    return tbl, shape


def _style_cell(
    cell,
    text,
    *,
    font_size=_PT_12,
    bold=False,
    color=AI_TEXT_LIGHT,
    alignment=PP_ALIGN.LEFT,
    font_name=FONT_BODY,
    fill_color=None,
):
    """Set text and style for a single table cell."""
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    _font(run, name=font_name, size=font_size, color=color, bold=bold)
    cell.text_frame.word_wrap = True
    if fill_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill_color
    # Reduce margins for compact look
    cell.margin_left = Inches(0.08)
    cell.margin_right = Inches(0.08)
    cell.margin_top = Inches(0.04)
    cell.margin_bottom = Inches(0.04)


# ---------------------------------------------------------------------------
# Slide functions
# ---------------------------------------------------------------------------


def slide_01_title(prs):
    """Slide 1 — Title slide."""
    slide = _blank_slide(prs)

    # Large accent bar near top
    add_accent_bar(slide, Inches(0), Inches(0.4), SLIDE_W, height=Pt(4))

    # Title
    add_textbox(
        slide,
        LEFT_MARGIN,
        Inches(2.2),
        CONTENT_W,
        Inches(1.2),
        text="ai-engineering",
        font_name=FONT_TITLE,
        font_size=Pt(52),
        color=AI_WHITE,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    # Subtitle line 1
    add_textbox(
        slide,
        LEFT_MARGIN,
        Inches(3.3),
        CONTENT_W,
        Inches(0.6),
        text="Governance for AI-Assisted Development",
        font_name=FONT_BODY,
        font_size=Pt(24),
        color=AI_TEXT_LIGHT,
        alignment=PP_ALIGN.CENTER,
    )

    # Accent bar below subtitle
    bar_w = Inches(3)
    bar_left = (SLIDE_W - bar_w) // 2
    add_accent_bar(slide, bar_left, Inches(4.05), bar_w, height=Pt(3))

    # Tagline
    add_textbox(
        slide,
        LEFT_MARGIN,
        Inches(4.3),
        CONTENT_W,
        Inches(0.5),
        text="Simple.  Efficient.  Practical.  Robust.  Secure.",
        font_name=FONT_BODY,
        font_size=Pt(16),
        color=AI_NEUTRAL,
        alignment=PP_ALIGN.CENTER,
    )

    # Badges
    badge_y = Inches(5.4)
    badge_w = Inches(1.6)
    badge_h = Inches(0.45)
    gap = Inches(0.3)
    badges = ["MIT License", "Python 3.11+", "Cross-IDE"]
    total_w = len(badges) * badge_w + (len(badges) - 1) * gap
    start_x = (SLIDE_W - total_w) // 2
    for i, label in enumerate(badges):
        x = start_x + i * (badge_w + gap)
        add_rounded_rect(
            slide,
            x,
            badge_y,
            badge_w,
            badge_h,
            fill_color=AI_CARD_DARK,
            border_color=AI_BORDER_DARK,
        )
        add_textbox(
            slide,
            x,
            badge_y,
            badge_w,
            badge_h,
            text=label,
            font_size=Pt(11),
            color=AI_NEUTRAL,
            alignment=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )

    # Bottom accent bar
    add_accent_bar(slide, Inches(0), Inches(7.0), SLIDE_W, height=Pt(4))

    set_notes(slide, NOTES[1])


def slide_02_evolution(prs):
    """Slide 2 — Evolution Timeline."""
    slide = _blank_slide(prs)
    add_slide_header(slide, "Current State of AI", subtitle="From autocomplete to governance")

    # Horizontal timeline line
    line_y = Inches(3.3)
    add_accent_bar(slide, LEFT_MARGIN, line_y, CONTENT_W, height=Pt(3))

    # Timeline nodes
    years = ["2022", "2023", "2024", "2025", "Today"]
    labels = [
        "Code\nCompletion",
        "Chat-in-\nIDE",
        "Agentic\nCoding",
        "Multi-Agent\n(MCP/A2A)",
        "Governance\n???",
    ]
    node_d = Inches(0.45)
    n = len(years)
    spacing = CONTENT_W / (n - 1) if n > 1 else 0
    for i in range(n):
        cx = LEFT_MARGIN + int(i * spacing) - node_d // 2
        is_last = i == n - 1
        fill = AI_ERROR if is_last else AI_ACCENT
        add_circle(slide, cx, line_y - node_d // 2 + Pt(1), node_d, fill_color=fill)
        # Year label above
        add_textbox(
            slide,
            cx - Inches(0.3),
            line_y - Inches(0.8),
            Inches(1.0),
            Inches(0.3),
            text=years[i],
            font_size=Pt(14),
            color=AI_TEXT_PRIMARY,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )
        # Description below
        add_textbox(
            slide,
            cx - Inches(0.4),
            line_y + Inches(0.4),
            Inches(1.2),
            Inches(0.6),
            text=labels[i],
            font_size=Pt(11),
            color=AI_TEXT_LIGHT,
            alignment=PP_ALIGN.CENTER,
        )

    # Definition cards at bottom
    defs = [
        ("Skills", "Reusable procedures\nin Markdown"),
        ("Agents", "Specialized personas —\nbehavior contracts"),
        ("MCP", "Model Context Protocol —\nexternal tool connection"),
        ("A2A", "Agent-to-Agent —\ncoordination between agents"),
    ]
    card_w = Inches(2.5)
    card_h = Inches(1.0)
    card_gap = Inches(0.2)
    total = len(defs) * card_w + (len(defs) - 1) * card_gap
    start_x = (SLIDE_W - total) // 2
    card_y = Inches(5.4)
    for i, (term, desc) in enumerate(defs):
        x = start_x + i * (card_w + card_gap)
        add_card(slide, x, card_y, card_w, card_h, fill_color=AI_CARD_DARK)
        add_textbox(
            slide,
            x + Inches(0.12),
            card_y + Inches(0.08),
            card_w - Inches(0.24),
            Inches(0.3),
            text=term,
            font_size=Pt(13),
            color=AI_ACCENT,
            bold=True,
        )
        add_textbox(
            slide,
            x + Inches(0.12),
            card_y + Inches(0.38),
            card_w - Inches(0.24),
            Inches(0.55),
            text=desc,
            font_size=Pt(10),
            color=AI_NEUTRAL,
        )

    set_notes(slide, NOTES[2])


def slide_03_problem(prs):
    """Slide 3 — The Problem: Ungoverned AI."""
    slide = _blank_slide(prs)
    add_slide_header(slide, "The Problem: Ungoverned AI")

    # Central codebase box
    cb_w, cb_h = Inches(2.2), Inches(1.0)
    cb_x = (SLIDE_W - cb_w) // 2
    cb_y = Inches(2.5)
    add_rect(
        slide,
        cb_x,
        cb_y,
        cb_w,
        cb_h,
        fill_color=AI_CARD_DARK,
        border_color=AI_BORDER_DARK,
        border_width=Pt(2),
    )
    add_textbox(
        slide,
        cb_x,
        cb_y,
        cb_w,
        cb_h,
        text="CODEBASE",
        font_size=Pt(16),
        color=AI_TEXT_PRIMARY,
        bold=True,
        alignment=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )

    # 4 agent boxes around the codebase
    agent_labels = ["Agent A", "Agent B", "Agent C", "Agent D"]
    agent_positions = [
        (cb_x - Inches(2.5), cb_y - Inches(0.1)),  # left
        (cb_x + cb_w + Inches(0.5), cb_y - Inches(0.1)),  # right
        (cb_x - Inches(1.3), cb_y - Inches(1.1)),  # top-left
        (cb_x + cb_w - Inches(0.9), cb_y - Inches(1.1)),  # top-right
    ]
    a_w, a_h = Inches(1.6), Inches(0.7)
    for label, (ax, ay) in zip(agent_labels, agent_positions, strict=True):
        add_rounded_rect(slide, ax, ay, a_w, a_h, fill_color=SEC_BLUE_PALE, border_color=SEC_BLUE)
        add_textbox(
            slide,
            ax,
            ay,
            a_w,
            a_h,
            text=label,
            font_size=Pt(12),
            color=SEC_BLUE_LIGHT,
            bold=True,
            alignment=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )

    # Warning triangles (small orange indicators between agents and codebase)
    tri_positions = [
        (cb_x - Inches(0.3), cb_y + Inches(0.3)),
        (cb_x + cb_w + Inches(0.05), cb_y + Inches(0.3)),
        (cb_x + Inches(0.3), cb_y - Inches(0.35)),
        (cb_x + cb_w - Inches(0.7), cb_y - Inches(0.35)),
    ]
    for tx, ty in tri_positions:
        # Use small orange diamond as warning indicator
        d = slide.shapes.add_shape(4, tx, ty, Inches(0.25), Inches(0.25))  # DIAMOND
        d.fill.solid()
        d.fill.fore_color.rgb = AI_ACCENT
        d.line.fill.background()

    # 4 risk cards at bottom
    risks = [
        ("Secrets in commits", "Credentials exposed\nin git history"),
        ("Quality gates bypassed", "No local enforcement;\ntests and lint ignored"),
        ("Architectural drift", "Each agent makes\ndifferent decisions"),
        ("Repeated decisions", "No memory between\nAI sessions"),
    ]
    card_w = Inches(2.5)
    card_h = Inches(1.15)
    card_gap = Inches(0.15)
    total = len(risks) * card_w + (len(risks) - 1) * card_gap
    start_x = (SLIDE_W - total) // 2
    card_y = Inches(4.3)
    for i, (title, desc) in enumerate(risks):
        x = start_x + i * (card_w + card_gap)
        add_card(
            slide,
            x,
            card_y,
            card_w,
            card_h,
            fill_color=AI_CARD_DARK,
            border_color=AI_BORDER_DARK,
            left_accent_color=AI_ACCENT,
        )
        add_textbox(
            slide,
            x + Inches(0.18),
            card_y + Inches(0.1),
            card_w - Inches(0.3),
            Inches(0.3),
            text=title,
            font_size=Pt(12),
            color=AI_TEXT_PRIMARY,
            bold=True,
        )
        add_textbox(
            slide,
            x + Inches(0.18),
            card_y + Inches(0.45),
            card_w - Inches(0.3),
            Inches(0.6),
            text=desc,
            font_size=Pt(10),
            color=AI_NEUTRAL,
        )

    # Bottom risk labels
    board_risks = ["Compliance gaps", "Security exposure", "Quality degradation", "Knowledge loss"]
    label_y = Inches(5.65)
    for i, r in enumerate(board_risks):
        x = start_x + i * (card_w + card_gap)
        add_textbox(
            slide,
            x + Inches(0.18),
            label_y,
            card_w - Inches(0.3),
            Inches(0.3),
            text=f"→ {r}",
            font_size=Pt(10),
            color=AI_ERROR,
            bold=True,
        )

    set_notes(slide, NOTES[3])


def slide_04_journey(prs):
    """Slide 4 — The Journey: 5 Rewrites."""
    slide = _blank_slide(prs)
    add_slide_header(slide, "The Journey: 5 Rewrites")

    # Input framework boxes (left column)
    frameworks = [
        ("SpecKit", "Specs without enforcement"),
        ("BMAD Method", "Multi-agent, heavyweight"),
        ("GSD", "Pragmatic, no governance"),
        ("OpenSpec", "Standard without tooling"),
    ]
    fw_x = LEFT_MARGIN
    fw_w = Inches(2.8)
    fw_h = Inches(0.7)
    fw_gap = Inches(0.2)
    fw_start_y = Inches(2.2)
    for i, (name, note) in enumerate(frameworks):
        y = fw_start_y + i * (fw_h + fw_gap)
        add_rect(slide, fw_x, y, fw_w, fw_h, fill_color=AI_CARD_DARK, border_color=AI_TEXT_MUTED)
        add_textbox(
            slide,
            fw_x + Inches(0.15),
            y + Inches(0.05),
            fw_w - Inches(0.3),
            Inches(0.3),
            text=name,
            font_size=Pt(14),
            color=AI_TEXT_PRIMARY,
            bold=True,
        )
        add_textbox(
            slide,
            fw_x + Inches(0.15),
            y + Inches(0.35),
            fw_w - Inches(0.3),
            Inches(0.3),
            text=note,
            font_size=Pt(10),
            color=AI_NEUTRAL,
        )

    # Chevron funnel in the middle
    chev_x = fw_x + fw_w + Inches(0.4)
    chev_y = Inches(3.0)
    add_chevron(slide, chev_x, chev_y, Inches(2.0), Inches(1.5), color=AI_ACCENT)
    add_textbox(
        slide,
        chev_x + Inches(0.2),
        chev_y + Inches(0.3),
        Inches(1.6),
        Inches(0.8),
        text="5\niterations",
        font_size=Pt(14),
        color=AI_BG_DARK,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )

    # Result box (right)
    res_x = chev_x + Inches(2.6)
    res_y = Inches(2.7)
    res_w = Inches(3.8)
    res_h = Inches(2.1)
    add_rect(
        slide,
        res_x,
        res_y,
        res_w,
        res_h,
        fill_color=AI_CARD_DARK,
        border_color=AI_ACCENT,
        border_width=Pt(2),
    )
    add_textbox(
        slide,
        res_x + Inches(0.2),
        res_y + Inches(0.15),
        res_w - Inches(0.4),
        Inches(0.4),
        text="ai-engineering",
        font_size=Pt(22),
        color=AI_TEXT_PRIMARY,
        bold=True,
    )
    add_accent_bar(slide, res_x + Inches(0.2), res_y + Inches(0.6), Inches(2.0))

    bullets = [
        "Content-first governance",
        "Simple to adopt",
        "Strict for enforcement",
        "Flexible to scale",
    ]
    for j, b in enumerate(bullets):
        add_textbox(
            slide,
            res_x + Inches(0.25),
            res_y + Inches(0.75) + j * Inches(0.3),
            res_w - Inches(0.5),
            Inches(0.3),
            text=f"→  {b}",
            font_size=Pt(12),
            color=AI_TEXT_LIGHT,
        )

    set_notes(slide, NOTES[4])


def slide_05_what_is(prs):
    """Slide 5 — What is ai-engineering."""
    slide = _blank_slide(prs)
    add_slide_header(
        slide,
        "What is ai-engineering?",
        subtitle="Content framework — not a platform, not a pipeline",
    )

    # Root header (orange-bordered box)
    root_x = LEFT_MARGIN
    root_y = Inches(2.2)
    root_w = Inches(3.0)
    root_h = Inches(0.55)
    add_rect(slide, root_x, root_y, root_w, root_h, fill_color=AI_ACCENT, border_color=None)
    add_textbox(
        slide,
        root_x,
        root_y,
        root_w,
        root_h,
        text=".ai-engineering/",
        font_size=Pt(16),
        color=AI_BG_DARK,
        bold=True,
        alignment=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )

    # 5 directory child cards
    dirs = [
        ("standards/", "Framework and\nteam rules"),
        ("skills/", "31 procedures\nin 7 categories"),
        ("agents/", "8 specialized\npersonas"),
        ("context/", "Specs, contracts,\nlearnings"),
        ("state/", "Decision store,\naudit log"),
    ]
    dir_w = Inches(2.0)
    dir_h = Inches(1.0)
    dir_gap = Inches(0.15)
    # Arrange in a row below the root
    dir_start_x = LEFT_MARGIN
    dir_y = Inches(3.1)
    for i, (name, desc) in enumerate(dirs):
        x = dir_start_x + i * (dir_w + dir_gap)
        add_card(slide, x, dir_y, dir_w, dir_h, fill_color=AI_CARD_DARK)
        add_textbox(
            slide,
            x + Inches(0.1),
            dir_y + Inches(0.08),
            dir_w - Inches(0.2),
            Inches(0.3),
            text=name,
            font_size=Pt(13),
            color=AI_ACCENT,
            bold=True,
        )
        add_textbox(
            slide,
            x + Inches(0.1),
            dir_y + Inches(0.4),
            dir_w - Inches(0.2),
            Inches(0.55),
            text=desc,
            font_size=Pt(10),
            color=AI_NEUTRAL,
        )

    # CLI bar
    cli_y = Inches(4.5)
    cli_w = Inches(6.5)
    cli_h = Inches(0.55)
    add_rect(slide, LEFT_MARGIN, cli_y, cli_w, cli_h, fill_color=AI_PRIMARY, border_color=None)
    add_textbox(
        slide,
        LEFT_MARGIN + Inches(0.2),
        cli_y,
        cli_w - Inches(0.4),
        cli_h,
        text="$ ai-eng install  |  update  |  doctor  |  validate",
        font_size=Pt(13),
        color=AI_TEXT_PRIMARY,
        bold=False,
        anchor=MSO_ANCHOR.MIDDLE,
    )

    # IDE cards (right side)
    ides = [
        ("Claude Code", "CLAUDE.md\n37 slash commands"),
        ("GitHub Copilot", "copilot-instructions.md\n.github/copilot/"),
        ("OpenAI Codex", "codex.md"),
    ]
    ide_x = LEFT_MARGIN + cli_w + Inches(0.3)
    ide_w = Inches(3.5)
    ide_h = Inches(0.65)
    ide_gap = Inches(0.15)
    ide_start_y = Inches(4.0)
    for i, (name, detail) in enumerate(ides):
        y = ide_start_y + i * (ide_h + ide_gap)
        add_card(
            slide, ide_x, y, ide_w, ide_h, fill_color=AI_CARD_DARK, border_color=AI_BORDER_DARK
        )
        add_textbox(
            slide,
            ide_x + Inches(0.12),
            y + Inches(0.05),
            Inches(1.5),
            Inches(0.3),
            text=name,
            font_size=Pt(12),
            color=AI_TEXT_PRIMARY,
            bold=True,
        )
        add_textbox(
            slide,
            ide_x + Inches(1.5),
            y + Inches(0.05),
            ide_w - Inches(1.6),
            Inches(0.55),
            text=detail,
            font_size=Pt(9),
            color=AI_NEUTRAL,
        )

    set_notes(slide, NOTES[5])


def slide_06_pipeline(prs):
    """Slide 6 — Developer Pipeline."""
    slide = _blank_slide(prs)
    add_slide_header(
        slide, "Developer Pipeline", subtitle="From install to governed commit in < 5 minutes"
    )

    # 5 stages as rounded rects with arrows
    stages = [
        ("ai-eng\ninstall .", SEC_TEAL),
        ("Code\n+ AI Assist", SEC_BLUE),
        ("Pre-\ncommit", AI_ACCENT),
        ("Commit-\nmsg", AI_PRIMARY),
        ("Pre-\npush", AI_ERROR),
    ]
    stage_w = Inches(1.6)
    stage_h = Inches(1.0)
    arrow_w = Inches(0.5)
    total_stages = len(stages) * stage_w + (len(stages) - 1) * arrow_w
    start_x = (SLIDE_W - total_stages) // 2
    stage_y = Inches(2.3)

    for i, (label, color) in enumerate(stages):
        x = start_x + i * (stage_w + arrow_w)
        add_rounded_rect(
            slide, x, stage_y, stage_w, stage_h, fill_color=AI_CARD_DARK, border_color=color
        )
        add_textbox(
            slide,
            x,
            stage_y,
            stage_w,
            stage_h,
            text=label,
            font_size=Pt(13),
            color=AI_TEXT_LIGHT,
            bold=True,
            alignment=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        # Arrow between stages
        if i < len(stages) - 1:
            ax = x + stage_w + Inches(0.05)
            add_arrow_right(
                slide, ax, stage_y + Inches(0.3), Inches(0.4), Inches(0.4), color=AI_ACCENT
            )

    # Detail cards below
    details = [
        ("Pre-commit", "ruff format + lint\ngitleaks secrets", "✓"),
        ("Commit-msg", "Valid format\nBranch protection", "✓"),
        ("Pre-push", "semgrep SAST/OWASP\npip-audit + pytest + ty", "✓"),
    ]
    card_w = Inches(3.2)
    card_h = Inches(1.3)
    card_gap = Inches(0.3)
    total_cards = len(details) * card_w + (len(details) - 1) * card_gap
    cards_start = (SLIDE_W - total_cards) // 2
    card_y = Inches(4.0)
    card_colors = [AI_ACCENT, AI_PRIMARY, AI_ERROR]

    for i, (title, desc, indicator) in enumerate(details):
        x = cards_start + i * (card_w + card_gap)
        add_card(
            slide,
            x,
            card_y,
            card_w,
            card_h,
            fill_color=AI_CARD_DARK,
            border_color=AI_BORDER_DARK,
            left_accent_color=card_colors[i],
        )
        add_textbox(
            slide,
            x + Inches(0.18),
            card_y + Inches(0.1),
            card_w - Inches(0.3),
            Inches(0.3),
            text=title,
            font_size=Pt(14),
            color=AI_TEXT_PRIMARY,
            bold=True,
        )
        add_textbox(
            slide,
            x + Inches(0.18),
            card_y + Inches(0.45),
            card_w - Inches(0.3),
            Inches(0.75),
            text=desc,
            font_size=Pt(11),
            color=AI_NEUTRAL,
        )
        # Green check indicator
        add_circle(
            slide, x + card_w - Inches(0.5), card_y + Inches(0.1), Inches(0.3), fill_color=SEC_TEAL
        )
        add_textbox(
            slide,
            x + card_w - Inches(0.5),
            card_y + Inches(0.1),
            Inches(0.3),
            Inches(0.3),
            text=indicator,
            font_size=Pt(14),
            color=AI_TEXT_PRIMARY,
            bold=True,
            alignment=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )

    # Bottom note
    add_textbox(
        slide,
        LEFT_MARGIN,
        Inches(5.7),
        CONTENT_W,
        Inches(0.4),
        text="If any gate fails → the push is blocked. No bypass. No --no-verify.",
        font_size=Pt(13),
        color=AI_ERROR,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )

    set_notes(slide, NOTES[6])


def slide_07_ownership(prs):
    """Slide 7 — Ownership Model."""
    slide = _blank_slide(prs)
    add_slide_header(slide, "Ownership Model", subtitle="Four boundaries, non-negotiable")

    layers = [
        (
            "Framework-managed",
            "standards/, skills/, agents/",
            "Updatable with ai-eng update",
            AI_ACCENT,
            "↻",
        ),
        ("Team-managed", "standards/team/", "NEVER overwritten by updates", SEC_BLUE, "🔒"),
        (
            "Project-managed",
            "context/, specs, learnings",
            "NEVER overwritten — institutional memory",
            SEC_TEAL,
            "🔒",
        ),
        ("System-managed", "state/ — runtime files", "Maintained automatically", SEC_PURPLE, "⚙"),
    ]
    layer_y = Inches(2.2)
    layer_h = Inches(1.05)
    layer_gap = Inches(0.15)
    layer_w = CONTENT_W

    for i, (name, path, desc, color, icon) in enumerate(layers):
        y = layer_y + i * (layer_h + layer_gap)
        # Main layer rect
        add_rect(
            slide,
            LEFT_MARGIN,
            y,
            layer_w,
            layer_h,
            fill_color=AI_CARD_DARK,
            border_color=color,
            border_width=Pt(2),
        )
        # Left accent
        add_rect(slide, LEFT_MARGIN, y, Inches(0.1), layer_h, fill_color=color, border_color=None)
        # Icon badge
        add_textbox(
            slide,
            LEFT_MARGIN + Inches(0.25),
            y + Inches(0.15),
            Inches(0.5),
            Inches(0.5),
            text=icon,
            font_size=Pt(22),
            color=color,
            alignment=PP_ALIGN.CENTER,
        )
        # Name
        add_textbox(
            slide,
            LEFT_MARGIN + Inches(0.8),
            y + Inches(0.1),
            Inches(3.0),
            Inches(0.4),
            text=name,
            font_size=Pt(18),
            color=AI_TEXT_PRIMARY,
            bold=True,
        )
        # Path
        add_textbox(
            slide,
            LEFT_MARGIN + Inches(0.8),
            y + Inches(0.55),
            Inches(3.5),
            Inches(0.35),
            text=path,
            font_size=Pt(12),
            color=AI_NEUTRAL,
        )
        # Description (right side)
        add_textbox(
            slide,
            LEFT_MARGIN + Inches(5.0),
            y + Inches(0.2),
            Inches(5.5),
            Inches(0.6),
            text=desc,
            font_size=Pt(13),
            color=AI_TEXT_LIGHT,
        )

    set_notes(slide, NOTES[7])


def slide_08_skills(prs):
    """Slide 8 — Skills Grid: 31 Procedures."""
    slide = _blank_slide(prs)
    add_slide_header(slide, "Skills: 31 Reusable Procedures")

    categories = [
        ("Workflows", "4", "commit, PR, acho,\npre-implementation", AI_ACCENT),
        ("Dev", "6", "debug, refactor, code-review,\ntest, migration, deps", SEC_BLUE),
        ("Review", "3", "architecture, performance,\nsecurity", SEC_TEAL),
        ("Docs", "4", "changelog, explain,\nwriter, prompt-design", SEC_PURPLE),
        ("Govern", "9", "create/delete specs, skills,\nagents + risk lifecycle", AI_PRIMARY),
        ("Quality", "3", "audit-code, audit-report,\ninstall-check", SEC_TEAL_DARK),
        ("Utils", "3", "git-helpers, platform-detect,\npython-patterns", SEC_BLUE_DARK),
    ]
    # 2 rows: 4 + 3
    card_w = Inches(2.5)
    card_h = Inches(1.8)
    card_gap = Inches(0.2)

    rows = [categories[:4], categories[4:]]
    row_y = Inches(2.1)

    for row_idx, row in enumerate(rows):
        total = len(row) * card_w + (len(row) - 1) * card_gap
        start_x = (SLIDE_W - total) // 2
        y = row_y + row_idx * (card_h + Inches(0.25))
        for i, (name, count, skills, color) in enumerate(row):
            x = start_x + i * (card_w + card_gap)
            add_card(
                slide,
                x,
                y,
                card_w,
                card_h,
                fill_color=AI_CARD_DARK,
                border_color=AI_BORDER_DARK,
                left_accent_color=color,
                left_accent_width=Inches(0.06),
            )
            # Category name
            add_textbox(
                slide,
                x + Inches(0.18),
                y + Inches(0.1),
                card_w - Inches(0.3),
                Inches(0.3),
                text=name,
                font_size=Pt(14),
                color=AI_TEXT_PRIMARY,
                bold=True,
            )
            # Count (large number)
            add_textbox(
                slide,
                x + Inches(0.18),
                y + Inches(0.4),
                Inches(0.6),
                Inches(0.6),
                text=count,
                font_size=Pt(32),
                color=color,
                bold=True,
            )
            # Skill list
            add_textbox(
                slide,
                x + Inches(0.8),
                y + Inches(0.5),
                card_w - Inches(1.0),
                Inches(0.5),
                text="skills",
                font_size=Pt(9),
                color=AI_TEXT_MUTED,
            )
            add_textbox(
                slide,
                x + Inches(0.18),
                y + Inches(1.05),
                card_w - Inches(0.3),
                Inches(0.65),
                text=skills,
                font_size=Pt(9),
                color=AI_NEUTRAL,
            )

    # Footer note
    add_textbox(
        slide,
        LEFT_MARGIN,
        Inches(6.5),
        CONTENT_W,
        Inches(0.3),
        text=(
            "They are NOT code — they are behavior specifications. "
            "Any AI agent reads and executes them."
        ),
        font_size=Pt(12),
        color=AI_NEUTRAL,
        alignment=PP_ALIGN.CENTER,
    )

    set_notes(slide, NOTES[8])


def slide_09_agents(prs):
    """Slide 9 — Agent Network."""
    slide = _blank_slide(prs)
    add_slide_header(slide, "Agents: 8 Specialized Personas")

    # Central hub
    hub_d = Inches(1.2)
    hub_x = (SLIDE_W - hub_d) // 2
    hub_y = Inches(3.6)
    add_circle(slide, hub_x, hub_y, hub_d, fill_color=AI_ACCENT)
    add_textbox(
        slide,
        hub_x,
        hub_y,
        hub_d,
        hub_d,
        text="ai-eng\nhub",
        font_size=Pt(12),
        color=AI_BG_DARK,
        bold=True,
        alignment=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )

    # 8 agent cards around the hub
    agents = [
        ("PE", "Principal\nEngineer", "Senior code review"),
        ("AR", "Architect", "Architectural analysis"),
        ("SR", "Security\nReviewer", "Security assessment"),
        ("DB", "Debugger", "Systematic diagnosis"),
        ("QA", "Quality\nAuditor", "Quality gates"),
        ("VA", "Verify\nApp", "E2E verification"),
        ("CM", "Codebase\nMapper", "Structure mapping"),
        ("CS", "Code\nSimplifier", "Complexity reduction"),
    ]
    # Positions around the hub (approximate circle layout)
    import math

    center_x = hub_x + hub_d // 2
    center_y = hub_y + hub_d // 2
    radius = Inches(2.4)
    agent_w = Inches(1.6)
    agent_h = Inches(1.1)
    init_d = Inches(0.4)

    for i, (initials, name, cap) in enumerate(agents):
        angle = math.pi / 2 + i * (2 * math.pi / len(agents))
        ax = int(center_x + radius * math.cos(angle)) - agent_w // 2
        ay = int(center_y - radius * math.sin(angle)) - agent_h // 2

        # Connector line (thin grey)
        line_shape = slide.shapes.add_shape(1, center_x, center_y, Inches(0.02), Inches(0.02))
        line_shape.fill.background()
        line_shape.line.fill.background()

        # Agent card
        add_card(
            slide, ax, ay, agent_w, agent_h, fill_color=AI_CARD_DARK, border_color=AI_BORDER_DARK
        )
        # Initials circle
        ix = ax + Inches(0.08)
        iy = ay + Inches(0.08)
        add_circle(slide, ix, iy, init_d, fill_color=AI_PRIMARY, border_color=AI_ACCENT)
        add_textbox(
            slide,
            ix,
            iy,
            init_d,
            init_d,
            text=initials,
            font_size=Pt(9),
            color=AI_ACCENT,
            bold=True,
            alignment=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        # Name
        add_textbox(
            slide,
            ax + Inches(0.5),
            ay + Inches(0.05),
            agent_w - Inches(0.6),
            Inches(0.5),
            text=name,
            font_size=Pt(10),
            color=AI_TEXT_PRIMARY,
            bold=True,
        )
        # Capability
        add_textbox(
            slide,
            ax + Inches(0.08),
            ay + Inches(0.65),
            agent_w - Inches(0.16),
            Inches(0.4),
            text=cap,
            font_size=Pt(8),
            color=AI_NEUTRAL,
            alignment=PP_ALIGN.CENTER,
        )

    set_notes(slide, NOTES[9])


def slide_10_spec_lifecycle(prs):
    """Slide 10 — Spec-Driven Delivery."""
    slide = _blank_slide(prs)
    add_slide_header(
        slide, "Spec-Driven Delivery", subtitle="4 documents, parallel execution, phase gates"
    )

    # 4 document boxes
    docs = [
        ("spec.md", "WHAT", "Requirements, scope,\nacceptance"),
        ("plan.md", "HOW", "Architecture,\ntrade-offs"),
        ("tasks.md", "DO", "Ordered tasks,\nassignable"),
        ("done.md", "DONE", "Completion,\nlearnings"),
    ]
    doc_w = Inches(2.2)
    doc_h = Inches(1.5)
    arrow_w = Inches(0.5)
    total = len(docs) * doc_w + (len(docs) - 1) * arrow_w
    start_x = (SLIDE_W - total) // 2
    doc_y = Inches(2.3)
    for i, (filename, phase, desc) in enumerate(docs):
        x = start_x + i * (doc_w + arrow_w)
        add_rect(
            slide,
            x,
            doc_y,
            doc_w,
            doc_h,
            fill_color=AI_CARD_DARK,
            border_color=AI_ACCENT,
            border_width=Pt(2),
        )
        # Phase label at top
        add_rect(slide, x, doc_y, doc_w, Inches(0.4), fill_color=AI_ACCENT, border_color=None)
        add_textbox(
            slide,
            x,
            doc_y,
            doc_w,
            Inches(0.4),
            text=phase,
            font_size=Pt(14),
            color=AI_BG_DARK,
            bold=True,
            alignment=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        # Filename
        add_textbox(
            slide,
            x + Inches(0.1),
            doc_y + Inches(0.5),
            doc_w - Inches(0.2),
            Inches(0.35),
            text=filename,
            font_size=Pt(14),
            color=AI_TEXT_PRIMARY,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )
        # Description
        add_textbox(
            slide,
            x + Inches(0.1),
            doc_y + Inches(0.9),
            doc_w - Inches(0.2),
            Inches(0.5),
            text=desc,
            font_size=Pt(11),
            color=AI_NEUTRAL,
            alignment=PP_ALIGN.CENTER,
        )
        # Arrow
        if i < len(docs) - 1:
            ax = x + doc_w + Inches(0.05)
            add_arrow_right(
                slide, ax, doc_y + Inches(0.55), Inches(0.4), Inches(0.4), color=AI_ACCENT
            )

    # Phase gate bar
    gate_y = Inches(4.2)
    gate_h = Inches(0.45)
    add_rect(slide, start_x, gate_y, total, gate_h, fill_color=AI_CARD_DARK, border_color=AI_ACCENT)
    add_textbox(
        slide,
        start_x,
        gate_y,
        total,
        gate_h,
        text="PHASE GATES  ·  Each phase passes a gate before the next",
        font_size=Pt(12),
        color=AI_TEXT_LIGHT,
        bold=True,
        alignment=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )

    # Branch lines showing parallel execution
    branch_y = Inches(5.1)
    branches = [
        ("Agent A", "feat/spec-NNN-phase1"),
        ("Agent B", "feat/spec-NNN-phase2"),
        ("Agent C", "feat/spec-NNN-phase3"),
    ]
    branch_w = Inches(3.2)
    branch_h = Inches(0.4)
    branch_gap = Inches(0.15)
    total_b = len(branches) * branch_w + (len(branches) - 1) * branch_gap
    b_start = (SLIDE_W - total_b) // 2
    for i, (agent, branch) in enumerate(branches):
        x = b_start + i * (branch_w + branch_gap)
        add_rect(
            slide,
            x,
            branch_y,
            branch_w,
            branch_h,
            fill_color=AI_CARD_DARK,
            border_color=AI_TEXT_MUTED,
        )
        add_textbox(
            slide,
            x + Inches(0.1),
            branch_y,
            Inches(1.0),
            branch_h,
            text=agent,
            font_size=Pt(10),
            color=AI_TEXT_PRIMARY,
            bold=True,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        add_textbox(
            slide,
            x + Inches(1.1),
            branch_y,
            branch_w - Inches(1.2),
            branch_h,
            text=branch,
            font_size=Pt(9),
            color=AI_NEUTRAL,
            anchor=MSO_ANCHOR.MIDDLE,
        )

    # Commit format
    add_textbox(
        slide,
        LEFT_MARGIN,
        Inches(5.8),
        CONTENT_W,
        Inches(0.35),
        text="Format: spec-NNN: Task X.Y — description  →  Each commit traceable to spec + task",
        font_size=Pt(12),
        color=AI_TEXT_LIGHT,
        alignment=PP_ALIGN.CENTER,
    )

    set_notes(slide, NOTES[10])


def slide_11_state(prs):
    """Slide 11 — State Management."""
    slide = _blank_slide(prs)
    add_slide_header(slide, "State Management and Decision Continuity")

    # 5 state file items
    state_files = [
        ("install-manifest.json", "What was installed, when, which version", "📦"),
        ("ownership-map.json", "Who owns each path", "🗺"),
        ("sources.lock.json", "Remote skills with verifiable checksums", "🔗"),
        ("decision-store.json", "10 real decisions with SHA-256 context hash", "💡"),
        ("audit-log.ndjson", "183 recorded events — append-only", "📋"),
    ]
    item_y = Inches(2.1)
    item_h = Inches(0.7)
    item_gap = Inches(0.12)
    item_w = Inches(6.0)

    for i, (filename, desc, icon) in enumerate(state_files):
        y = item_y + i * (item_h + item_gap)
        # Icon
        add_textbox(
            slide,
            LEFT_MARGIN,
            y,
            Inches(0.5),
            item_h,
            text=icon,
            font_size=Pt(22),
            alignment=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
            color=AI_ACCENT,
        )
        # Filename
        add_textbox(
            slide,
            LEFT_MARGIN + Inches(0.6),
            y + Inches(0.05),
            Inches(2.8),
            Inches(0.3),
            text=filename,
            font_size=Pt(14),
            color=AI_TEXT_PRIMARY,
            bold=True,
        )
        # Description
        add_textbox(
            slide,
            LEFT_MARGIN + Inches(0.6),
            y + Inches(0.35),
            item_w - Inches(0.6),
            Inches(0.3),
            text=desc,
            font_size=Pt(12),
            color=AI_NEUTRAL,
        )

    # Decision continuity box (right side)
    dc_x = LEFT_MARGIN + item_w + Inches(0.5)
    dc_y = Inches(2.1)
    dc_w = Inches(4.5)
    dc_h = Inches(4.2)
    add_rect(
        slide,
        dc_x,
        dc_y,
        dc_w,
        dc_h,
        fill_color=AI_CARD_DARK,
        border_color=AI_ACCENT,
        border_width=Pt(2),
    )
    add_textbox(
        slide,
        dc_x + Inches(0.2),
        dc_y + Inches(0.15),
        dc_w - Inches(0.4),
        Inches(0.35),
        text="Decision Continuity",
        font_size=Pt(16),
        color=AI_TEXT_PRIMARY,
        bold=True,
    )
    add_accent_bar(slide, dc_x + Inches(0.2), dc_y + Inches(0.55), Inches(2.5))

    dc_text = (
        "Agent A decides in session 1\n"
        "  → generates SHA-256 context hash\n\n"
        "Agent B arrives in session 5\n"
        "  → reads decision store\n"
        "  → does NOT ask again\n\n"
        "Only re-prompt if:\n"
        "  • expired\n"
        "  • scope changed\n"
        "  • severity changed\n"
        "  • policy changed\n"
        "  • context hash changed"
    )
    add_textbox(
        slide,
        dc_x + Inches(0.2),
        dc_y + Inches(0.75),
        dc_w - Inches(0.4),
        Inches(3.2),
        text=dc_text,
        font_size=Pt(11),
        color=AI_TEXT_LIGHT,
    )

    set_notes(slide, NOTES[11])


def slide_12_quality_gates(prs):
    """Slide 12 — Quality Gates & Security."""
    slide = _blank_slide(prs)
    add_slide_header(slide, "Quality Gates and Security")

    # 3 stage columns
    stages = [
        ("Pre-commit", ["ruff format", "ruff lint", "gitleaks"], AI_ACCENT),
        ("Commit-msg", ["Valid format", "Branch protection"], AI_PRIMARY),
        ("Pre-push", ["semgrep SAST", "pip-audit CVE", "pytest", "ty types"], AI_ERROR),
    ]
    col_w = Inches(3.2)
    col_gap = Inches(0.3)
    total = len(stages) * col_w + (len(stages) - 1) * col_gap
    start_x = (SLIDE_W - total) // 2
    col_y = Inches(2.1)

    for i, (name, checks, color) in enumerate(stages):
        x = start_x + i * (col_w + col_gap)
        # Header
        add_rect(slide, x, col_y, col_w, Inches(0.45), fill_color=color, border_color=None)
        add_textbox(
            slide,
            x,
            col_y,
            col_w,
            Inches(0.45),
            text=name,
            font_size=Pt(14),
            color=AI_BG_DARK,
            bold=True,
            alignment=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        # Check items
        for j, check in enumerate(checks):
            cy = col_y + Inches(0.55) + j * Inches(0.35)
            add_textbox(
                slide,
                x + Inches(0.15),
                cy,
                col_w - Inches(0.3),
                Inches(0.3),
                text=f"✓  {check}",
                font_size=Pt(12),
                color=AI_TEXT_LIGHT,
            )

    # Threshold metric cards
    metrics = [
        ("Coverage", "≥ 80%", "(≥90% gov-critical)"),
        ("Duplication", "≤ 3%", ""),
        ("Cyclomatic", "≤ 10", "complexity"),
        ("Cognitive", "≤ 15", "complexity"),
    ]
    metric_w = Inches(2.3)
    metric_h = Inches(0.9)
    metric_gap = Inches(0.2)
    total_m = len(metrics) * metric_w + (len(metrics) - 1) * metric_gap
    m_start = (SLIDE_W - total_m) // 2
    m_y = Inches(4.3)

    for i, (label, value, extra) in enumerate(metrics):
        x = m_start + i * (metric_w + metric_gap)
        add_card(slide, x, m_y, metric_w, metric_h, fill_color=AI_CARD_DARK)
        add_textbox(
            slide,
            x + Inches(0.1),
            m_y + Inches(0.08),
            metric_w - Inches(0.2),
            Inches(0.25),
            text=label,
            font_size=Pt(11),
            color=AI_NEUTRAL,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )
        add_textbox(
            slide,
            x + Inches(0.1),
            m_y + Inches(0.3),
            metric_w - Inches(0.2),
            Inches(0.35),
            text=value,
            font_size=Pt(20),
            color=AI_ACCENT,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )
        if extra:
            add_textbox(
                slide,
                x + Inches(0.1),
                m_y + Inches(0.62),
                metric_w - Inches(0.2),
                Inches(0.2),
                text=extra,
                font_size=Pt(8),
                color=AI_TEXT_MUTED,
                alignment=PP_ALIGN.CENTER,
            )

    # Risk severity cards
    severities = [
        ("Critical", "15 days", AI_ERROR),
        ("High", "30 days", AI_PRIMARY),
        ("Medium", "60 days", AI_ACCENT),
        ("Low", "90 days", AI_TEXT_MUTED),
    ]
    sev_w = Inches(2.3)
    sev_h = Inches(0.65)
    sev_gap = Inches(0.2)
    total_s = len(severities) * sev_w + (len(severities) - 1) * sev_gap
    s_start = (SLIDE_W - total_s) // 2
    s_y = Inches(5.6)

    for i, (level, expiry, color) in enumerate(severities):
        x = s_start + i * (sev_w + sev_gap)
        add_card(
            slide,
            x,
            s_y,
            sev_w,
            sev_h,
            fill_color=AI_CARD_DARK,
            border_color=AI_BORDER_DARK,
            left_accent_color=color,
        )
        add_textbox(
            slide,
            x + Inches(0.18),
            s_y + Inches(0.05),
            Inches(1.0),
            Inches(0.25),
            text=level,
            font_size=Pt(12),
            color=AI_TEXT_PRIMARY,
            bold=True,
        )
        add_textbox(
            slide,
            x + Inches(0.18),
            s_y + Inches(0.32),
            sev_w - Inches(0.3),
            Inches(0.25),
            text=f"Expiry: {expiry}  ·  Max 2 renewals",
            font_size=Pt(9),
            color=AI_NEUTRAL,
        )

    set_notes(slide, NOTES[12])


def slide_13_value_by_role(prs):
    """Slide 13 — Value by Role (table)."""
    slide = _blank_slide(prs)
    add_slide_header(slide, "Value by Role")

    # 5-row table: Role / Value / Key Feature
    data = [
        ("Role", "Value", "Key Feature"),
        ("Engineers", "Fast workflows — 37 slash commands", "Quality gates before push"),
        ("Governance", "audit-log with 183 traceable events", "Decision store with 10 decisions"),
        (
            "Security / AppSec",
            "gitleaks + semgrep + pip-audit on every push",
            "Risk acceptance with expiry",
        ),
        ("Quality / DevEx", "Sonar-like quality gates WITHOUT a server", "Setup in < 5 minutes"),
        (
            "Architecture",
            "Standards with layering and ownership",
            "Agent personas for consistent review",
        ),
    ]

    tbl_left = LEFT_MARGIN
    tbl_top = Inches(2.1)
    tbl_w = CONTENT_W
    tbl_h = Inches(4.0)
    rows = len(data)
    cols = 3

    tbl, _shape = add_styled_table(slide, tbl_left, tbl_top, tbl_w, tbl_h, rows, cols)

    # Set column widths
    tbl.columns[0].width = Inches(2.2)
    tbl.columns[1].width = Inches(5.0)
    tbl.columns[2].width = CONTENT_W - Inches(7.2)

    for ri, row_data in enumerate(data):
        for ci, cell_text in enumerate(row_data):
            cell = tbl.cell(ri, ci)
            if ri == 0:
                _style_cell(
                    cell,
                    cell_text,
                    font_size=Pt(13),
                    bold=True,
                    color=AI_BG_DARK,
                    font_name=FONT_TITLE,
                    fill_color=AI_ACCENT,
                )
            else:
                fill = AI_CARD_DARK if ri % 2 == 0 else None
                _style_cell(
                    cell,
                    cell_text,
                    font_size=Pt(11),
                    bold=(ci == 0),
                    color=AI_TEXT_PRIMARY if ci == 0 else AI_TEXT_LIGHT,
                    fill_color=fill,
                )

    set_notes(slide, NOTES[13])


def slide_14_business_case(prs):
    """Slide 14 — Business Case: Metric Cards."""
    slide = _blank_slide(prs)
    add_slide_header(slide, "Business Case: ROI and Risk Reduction")

    metrics = [
        ("100%", "Gate Execution", "0 ungated operations.\nFull local enforcement."),
        ("$0", "License Cost", "MIT open source.\nNo SonarQube, no additional CI."),
        ("< 5 min", "Time to Governed", "From install to\nfirst governed commit."),
        ("SHA-256", "Compliance Ready", "Audit log + risk acceptance\n+ decision store."),
    ]
    card_w = Inches(2.5)
    card_h = Inches(3.0)
    card_gap = Inches(0.2)
    total = len(metrics) * card_w + (len(metrics) - 1) * card_gap
    start_x = (SLIDE_W - total) // 2
    card_y = Inches(2.3)

    for i, (number, label, detail) in enumerate(metrics):
        x = start_x + i * (card_w + card_gap)
        add_card(
            slide, x, card_y, card_w, card_h, fill_color=AI_CARD_DARK, border_color=AI_BORDER_DARK
        )

        # Big number
        add_textbox(
            slide,
            x + Inches(0.1),
            card_y + Inches(0.3),
            card_w - Inches(0.2),
            Inches(0.8),
            text=number,
            font_size=Pt(36),
            color=AI_ACCENT,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )
        # Label
        add_textbox(
            slide,
            x + Inches(0.1),
            card_y + Inches(1.1),
            card_w - Inches(0.2),
            Inches(0.4),
            text=label,
            font_size=Pt(15),
            color=AI_TEXT_PRIMARY,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )
        # Accent bar
        add_accent_bar(slide, x + Inches(0.4), card_y + Inches(1.55), card_w - Inches(0.8))
        # Detail
        add_textbox(
            slide,
            x + Inches(0.15),
            card_y + Inches(1.75),
            card_w - Inches(0.3),
            Inches(1.0),
            text=detail,
            font_size=Pt(11),
            color=AI_NEUTRAL,
            alignment=PP_ALIGN.CENTER,
        )

    set_notes(slide, NOTES[14])


def slide_15_multi_ide(prs):
    """Slide 15 — Multi-IDE Support."""
    slide = _blank_slide(prs)
    add_slide_header(
        slide, "Multi-IDE: One Framework, Every Provider", subtitle="No vendor lock-in"
    )

    # 4 IDE boxes converging to central .ai-engineering/
    ides = [
        ("Claude Code", "CLAUDE.md\n37 slash commands\n.claude/commands/"),
        ("GitHub Copilot", "copilot-instructions.md\n.github/copilot/"),
        ("OpenAI Codex", "codex.md"),
        ("Terminal CLI", "ai-eng install\nai-eng doctor"),
    ]

    # Central box
    center_w = Inches(2.8)
    center_h = Inches(1.2)
    center_x = (SLIDE_W - center_w) // 2
    center_y = Inches(3.5)
    add_rect(
        slide,
        center_x,
        center_y,
        center_w,
        center_h,
        fill_color=AI_CARD_DARK,
        border_color=AI_ACCENT,
        border_width=Pt(3),
    )
    add_textbox(
        slide,
        center_x,
        center_y,
        center_w,
        center_h,
        text=".ai-engineering/",
        font_size=Pt(18),
        color=AI_TEXT_PRIMARY,
        bold=True,
        alignment=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )

    # IDE boxes: 2 on left, 2 on right
    ide_w = Inches(2.8)
    ide_h = Inches(1.1)
    positions = [
        (LEFT_MARGIN, Inches(2.2)),
        (LEFT_MARGIN, Inches(4.6)),
        (SLIDE_W - RIGHT_MARGIN - ide_w, Inches(2.2)),
        (SLIDE_W - RIGHT_MARGIN - ide_w, Inches(4.6)),
    ]
    for i, ((ix, iy), (name, detail)) in enumerate(zip(positions, ides, strict=True)):
        add_rect(slide, ix, iy, ide_w, ide_h, fill_color=AI_CARD_DARK, border_color=AI_TEXT_MUTED)
        add_textbox(
            slide,
            ix + Inches(0.15),
            iy + Inches(0.08),
            ide_w - Inches(0.3),
            Inches(0.3),
            text=name,
            font_size=Pt(14),
            color=AI_TEXT_PRIMARY,
            bold=True,
        )
        add_textbox(
            slide,
            ix + Inches(0.15),
            iy + Inches(0.4),
            ide_w - Inches(0.3),
            Inches(0.65),
            text=detail,
            font_size=Pt(10),
            color=AI_NEUTRAL,
        )

        # Arrows pointing to center
        if i < 2:  # left side
            ax = ix + ide_w + Inches(0.1)
            ay = iy + ide_h // 2 - Inches(0.15)
            add_arrow_right(
                slide, ax, ay, center_x - ax - Inches(0.1), Inches(0.3), color=AI_ACCENT
            )
        else:  # right side — arrow from center to right (reversed visually)
            ax = center_x + center_w + Inches(0.1)
            ay = iy + ide_h // 2 - Inches(0.15)
            target_x = ix - Inches(0.1)
            # Use a left arrow via flipping — just use a bar + triangle
            add_accent_bar(
                slide, ax, ay + Inches(0.12), target_x - ax, height=Pt(3), color=AI_ACCENT
            )

    set_notes(slide, NOTES[15])


def slide_16_comparison(prs):
    """Slide 16 — Plain AI vs ai-engineering comparison table."""
    slide = _blank_slide(prs)
    add_slide_header(slide, "Why not AI with simple instructions?")

    capabilities = [
        ("Capability", "Plain AI", "ai-engineering"),
        ("Local enforcement", "—", "✓"),
        ("State / memory", "—", "✓"),
        ("Ownership model", "—", "✓"),
        ("Audit trail", "—", "✓"),
        ("Security scanning", "—", "✓"),
        ("Risk management", "—", "✓"),
        ("Delivery lifecycle", "—", "✓"),
        ("Cross-IDE governance", "—", "✓"),
    ]

    tbl_left = Inches(2.5)
    tbl_top = Inches(2.1)
    tbl_w = Inches(8.0)
    tbl_h = Inches(4.2)
    rows = len(capabilities)
    cols = 3

    tbl, _shape = add_styled_table(slide, tbl_left, tbl_top, tbl_w, tbl_h, rows, cols)

    tbl.columns[0].width = Inches(3.2)
    tbl.columns[1].width = Inches(2.0)
    tbl.columns[2].width = Inches(2.8)

    for ri, row_data in enumerate(capabilities):
        for ci, cell_text in enumerate(row_data):
            cell = tbl.cell(ri, ci)
            if ri == 0:
                _style_cell(
                    cell,
                    cell_text,
                    font_size=Pt(13),
                    bold=True,
                    color=AI_BG_DARK,
                    font_name=FONT_TITLE,
                    fill_color=AI_ACCENT,
                    alignment=PP_ALIGN.CENTER,
                )
            else:
                is_check = cell_text == "✓"
                is_dash = cell_text == "—"
                if is_check:
                    _style_cell(
                        cell,
                        cell_text,
                        font_size=Pt(16),
                        bold=True,
                        color=SEC_TEAL,
                        alignment=PP_ALIGN.CENTER,
                        fill_color=AI_CARD_DARK if ri % 2 == 0 else None,
                    )
                elif is_dash:
                    _style_cell(
                        cell,
                        cell_text,
                        font_size=Pt(16),
                        bold=True,
                        color=AI_TEXT_MUTED,
                        alignment=PP_ALIGN.CENTER,
                        fill_color=AI_CARD_DARK if ri % 2 == 0 else None,
                    )
                else:
                    _style_cell(
                        cell,
                        cell_text,
                        font_size=Pt(11),
                        bold=(ci == 0),
                        color=AI_TEXT_LIGHT,
                        fill_color=AI_CARD_DARK if ri % 2 == 0 else None,
                    )

    # Score bar
    add_textbox(
        slide,
        Inches(2.5),
        Inches(6.5),
        Inches(8.0),
        Inches(0.4),
        text=(
            "Plain AI: 0/8 capabilities enforced    ·    ai-engineering: 8/8 capabilities enforced"
        ),
        font_size=Pt(14),
        color=AI_ERROR,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )

    set_notes(slide, NOTES[16])


def slide_17_frameworks(prs):
    """Slide 17 — Frameworks Comparison (scored table)."""
    slide = _blank_slide(prs)
    add_slide_header(slide, "Comparison with Alternatives", subtitle="Scored comparison — 6 axes")

    # Data: axis / SpecKit / BMAD / GSD / OpenSpec / ai-engineering
    axes = [
        ("Axis", "SpecKit", "BMAD", "GSD", "OpenSpec", "ai-eng"),
        ("Governance depth", "3", "5", "2", "4", "9"),
        ("Enforcement", "1", "3", "1", "2", "9"),
        ("Practicality", "6", "3", "8", "2", "8"),
        ("Cross-IDE", "3", "2", "4", "5", "9"),
        ("Risk lifecycle", "0", "1", "0", "2", "8"),
        ("State continuity", "2", "3", "1", "3", "9"),
    ]

    tbl_left = LEFT_MARGIN
    tbl_top = Inches(2.2)
    tbl_w = CONTENT_W
    tbl_h = Inches(3.8)
    rows = len(axes)
    cols = 6

    tbl, _shape = add_styled_table(slide, tbl_left, tbl_top, tbl_w, tbl_h, rows, cols)

    tbl.columns[0].width = Inches(2.2)
    remaining = CONTENT_W - Inches(2.2)
    for ci in range(1, cols):
        tbl.columns[ci].width = int(remaining / 5)

    # Color code scores
    def _score_color(score_str):
        try:
            s = int(score_str)
        except ValueError:
            return None
        if s >= 8:
            return SEC_TEAL_PALE
        if s >= 5:
            return SEC_BLUE_PALE
        if s >= 3:
            return AI_CARD_DARK
        return None  # low scores get no fill

    for ri, row_data in enumerate(axes):
        for ci, cell_text in enumerate(row_data):
            cell = tbl.cell(ri, ci)
            if ri == 0:
                # Last column header highlighted
                fill = AI_ACCENT if ci == 5 else AI_PRIMARY
                text_c = AI_BG_DARK if ci == 5 else AI_TEXT_PRIMARY
                _style_cell(
                    cell,
                    cell_text,
                    font_size=Pt(12),
                    bold=True,
                    color=text_c,
                    font_name=FONT_TITLE,
                    fill_color=fill,
                    alignment=PP_ALIGN.CENTER,
                )
            else:
                is_axis = ci == 0
                is_aieng = ci == 5
                fill = _score_color(cell_text) if ci > 0 else None
                _style_cell(
                    cell,
                    cell_text,
                    font_size=Pt(13) if not is_axis else Pt(11),
                    bold=is_axis or is_aieng,
                    color=AI_ACCENT
                    if is_aieng
                    else (AI_TEXT_LIGHT if is_axis else AI_TEXT_PRIMARY),
                    alignment=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT,
                    fill_color=fill,
                )

    # Legend
    add_textbox(
        slide,
        LEFT_MARGIN,
        Inches(6.2),
        CONTENT_W,
        Inches(0.4),
        text=(
            "Scale 0-10  ·  Green >=8 (leader)  ·  Blue >=5 (competent)"
            "  ·  Grey >=3 (partial)  ·  No color <3"
        ),
        font_size=Pt(11),
        color=AI_NEUTRAL,
        alignment=PP_ALIGN.CENTER,
    )

    # Key differentiators
    diffs = (
        "5 differentiators: Content-first  ·  Non-bypassable enforcement"
        "  ·  Risk lifecycle  ·  Cross-IDE day one  ·  Ownership model"
    )
    add_textbox(
        slide,
        LEFT_MARGIN,
        Inches(6.6),
        CONTENT_W,
        Inches(0.4),
        text=diffs,
        font_size=Pt(11),
        color=AI_TEXT_LIGHT,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )

    set_notes(slide, NOTES[17])


def slide_18_cta(prs):
    """Slide 18 — Call to Action."""
    slide = _blank_slide(prs)
    add_slide_header(slide, "")

    # "The Ask" in large orange
    add_textbox(
        slide,
        LEFT_MARGIN,
        Inches(1.2),
        CONTENT_W,
        Inches(0.8),
        text="The Ask",
        font_name=FONT_TITLE,
        font_size=Pt(44),
        color=AI_ACCENT,
        bold=True,
        alignment=PP_ALIGN.LEFT,
    )

    add_textbox(
        slide,
        LEFT_MARGIN,
        Inches(2.0),
        CONTENT_W,
        Inches(0.5),
        text="Approve ai-engineering as the governance standard for AI-assisted development.",
        font_size=Pt(18),
        color=AI_TEXT_PRIMARY,
        bold=True,
    )

    # Pilot plan bullets
    pilot_bullets = [
        "→  2-3 repositories next quarter",
        "→  Metrics: gate execution, time to governed commit,",
        "     security catch rate, decision reuse rate",
    ]
    for j, b in enumerate(pilot_bullets):
        add_textbox(
            slide,
            LEFT_MARGIN + Inches(0.2),
            Inches(2.7) + j * Inches(0.32),
            Inches(8.0),
            Inches(0.3),
            text=b,
            font_size=Pt(14),
            color=AI_TEXT_LIGHT,
        )

    # 3-phase horizontal roadmap bars
    phases = [
        ("Phase 1 — Now", "GitHub + Python + Claude/Copilot/Codex", AI_ACCENT),
        ("Phase 2", "Azure DevOps + more stacks + signature verification", AI_PRIMARY),
        ("Phase 3", "Multi-agent orchestration + docs site", AI_ERROR),
    ]
    phase_y = Inches(4.2)
    phase_h = Inches(0.55)
    phase_gap = Inches(0.15)
    phase_widths = [Inches(4.5), Inches(3.5), Inches(2.8)]

    for i, ((label, desc, color), pw) in enumerate(zip(phases, phase_widths, strict=True)):
        y = phase_y + i * (phase_h + phase_gap)
        add_rect(slide, LEFT_MARGIN, y, pw, phase_h, fill_color=color, border_color=None)
        add_textbox(
            slide,
            LEFT_MARGIN + Inches(0.15),
            y,
            pw - Inches(0.3),
            phase_h,
            text=label,
            font_size=Pt(13),
            color=AI_BG_DARK,
            bold=True,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        # Description next to bar
        add_textbox(
            slide,
            LEFT_MARGIN + pw + Inches(0.2),
            y,
            Inches(6.0),
            phase_h,
            text=desc,
            font_size=Pt(12),
            color=AI_TEXT_LIGHT,
            anchor=MSO_ANCHOR.MIDDLE,
        )

    # Investment summary
    add_rect(
        slide,
        LEFT_MARGIN,
        Inches(6.1),
        CONTENT_W,
        Inches(0.5),
        fill_color=AI_CARD_DARK,
        border_color=None,
    )
    add_textbox(
        slide,
        LEFT_MARGIN,
        Inches(6.1),
        CONTENT_W,
        Inches(0.5),
        text=(
            "Investment: $0 licenses (MIT)  ·  Tooling: Python + git hooks"
            "  ·  Framework ready to install today."
        ),
        font_size=Pt(13),
        color=AI_TEXT_LIGHT,
        bold=True,
        alignment=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )

    # Bottom accent bar
    add_accent_bar(slide, Inches(0), Inches(7.0), SLIDE_W, height=Pt(4))

    set_notes(slide, NOTES[18])


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_01_title(prs)
    slide_02_evolution(prs)
    slide_03_problem(prs)
    slide_04_journey(prs)
    slide_05_what_is(prs)
    slide_06_pipeline(prs)
    slide_07_ownership(prs)
    slide_08_skills(prs)
    slide_09_agents(prs)
    slide_10_spec_lifecycle(prs)
    slide_11_state(prs)
    slide_12_quality_gates(prs)
    slide_13_value_by_role(prs)
    slide_14_business_case(prs)
    slide_15_multi_ide(prs)
    slide_16_comparison(prs)
    slide_17_frameworks(prs)
    slide_18_cta(prs)

    out = Path(__file__).parent / "ai-engineering-board.pptx"
    prs.save(str(out))
    print(f"Generated: {out}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
